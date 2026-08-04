"""Deck v2 = the existing 12-slide method deck + the budget-conditioning section.

Reads `method_deck_v1.pptx` (never modified) and writes a NEW file
`method_deck_v2.pptx` into this experiment directory: v1's slides in order,
the new section inserted before the closing slide, and a refreshed closing
slide reflecting the budget results.
"""

import copy
import json
import os

from pptx import Presentation
from pptx.util import Emu, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
V1 = "/home/yilin/Research/Route_TSC_CART/experiments/20260716-fm-mcvrp-local/ppt/method_deck_v1.pptx"
OUT = f"{HERE}/ppt/method_deck_v2.pptx"
R = f"{HERE}/results"
EMU_IN = 914400

AGG = json.load(open(f"{R}/agg_3seed.json"))
CURVE = json.load(open(f"{R}/curve.json"))


def pct(k):
    return (f"{100*AGG[k]['rel_gap_mean']:.1f}% ± {100*AGG[k]['rel_gap_std']:.1f}")


def clone_content(prs, src):
    """Clone a content slide, keeping only its placeholders (header/subtitle)."""
    new = prs.slides.add_slide(src.slide_layout)
    for shp in list(new.shapes):
        shp._element.getparent().remove(shp._element)
    for shp in src.shapes:
        if shp.is_placeholder:
            new.shapes._spTree.append(copy.deepcopy(shp._element))
    return new


def set_ph(slide, idx, text):
    for sh in slide.shapes:
        if sh.is_placeholder and sh.placeholder_format.idx == idx:
            p = sh.text_frame.paragraphs[0]
            if p.runs:
                p.runs[0].text = text
                for r in p.runs[1:]:
                    r._r.getparent().remove(r._r)
            else:
                p.text = text
            return sh
    return None


def add_bullets(slide, x_in, y_in, w_in, lines, size=14):
    tb = slide.shapes.add_textbox(Emu(int(x_in * EMU_IN)), Emu(int(y_in * EMU_IN)),
                                  Emu(int(w_in * EMU_IN)), Emu(EMU_IN))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (lvl, txt, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " if lvl == 0 else "   – ") + txt if lvl >= 0 else txt
        p.font.size = Pt(size if lvl <= 0 else size - 1.5)
        p.font.bold = bold
    return tb


def add_pic(slide, path, x_in, y_in, w_in):
    return slide.shapes.add_picture(path, Emu(int(x_in * EMU_IN)),
                                    Emu(int(y_in * EMU_IN)),
                                    width=Emu(int(w_in * EMU_IN)))


prs = Presentation(V1)
S = list(prs.slides)
n_v1 = len(S)
content_src = S[2]                      # a "Methodology / Problem" content slide
closing_src = S[-1]                     # "Summary / Next steps"
print(f"v1 has {n_v1} slides")

new = []

# ---------------------------------------------------------------- S13
s = clone_content(prs, content_src)
set_ph(s, 11, "Extension 2")
set_ph(s, 10, "The weight α is not a dial")
add_bullets(s, 0.7, 1.55, 11.9, [
    (0, "Objective:   min [ α₁·Σcost − α₂·Σ y(i,τ) ] / (V·H),   both terms verified in (0,1] "
        "before α (max exactly 1.000 across all 6 100 labels — normalization is NOT the issue)", False),
    (0, "Marginal argument: one extra step through a fresh cell costs ΔC = 1 and buys ΔK ≤ 1, so", True),
    (-1, "        Δobj = (α₁ − α₂) / (V·H)      →  sign flips at α₂/α₁ = 1, independent of any normalizer", True),
    (0, "α₂ > α₁ ⇒ EVERY detour is profitable ⇒ the optimum roams until the horizon stops it; "
        "α₂ < α₁ ⇒ no detour is ever bought (coverage survives only as a tie-break)", False),
    (0, "So the real limit on coverage is the HORIZON H, not α. α selects one of two regimes; "
        "it cannot set a position between them", True),
    (0, "Consequence for learning: conditioning a neural solver on α would teach it a switch, "
        "not a continuous semantics — not worth new labels", False),
], size=14.5)

# ---------------------------------------------------------------- S14
s = clone_content(prs, content_src)
set_ph(s, 11, "Extension 2")
set_ph(s, 10, "Switch vs dial — measured")
add_pic(s, f"{R}/figS1_knob_vs_switch.png", 0.77, 1.5, 11.8)
add_bullets(s, 0.7, 6.10, 11.9, [
    (0, "One fixed instance, one MILP, one solver, MIPGap 0.5%. LEFT: 13 values of α₂, budget fixed at H "
        "→ 2 distinct solutions (α₂ = 0.10–0.48 identical link-for-link). RIGHT: α frozen at 0.3/0.7, "
        "budget swept → 10 distinct monotone solutions", False),
], size=12.5)

# ---------------------------------------------------------------- S15
s = clone_content(prs, content_src)
set_ph(s, 11, "Extension 2")
set_ph(s, 10, "The fix: a per-vehicle budget")
add_bullets(s, 0.7, 1.62, 11.9, [
    (0, "Every vehicle carries its own budget Bᵥ (= shift length / remaining energy — a PHYSICAL "
        "parameter, not a tuning weight); hard deadline  t0ᵥ + Bᵥ.  The objective is unchanged:", True),
    (-1, "        min [ α₁·Σcost − α₂·Σ y(i,τ) ] / Σᵥ Bᵥ        s.t.  vehicle v reaches dᵥ by t0ᵥ + Bᵥ", True),
    (0, "Parameterised as a SLACK RATIO:  Bᵥ = ⌈ρᵥ · τᵐⁱⁿ(oᵥ,dᵥ,t0ᵥ)⌉ — an absolute budget is not "
        "comparable across ODs (vacuous for near pairs, infeasible for far ones). ρ is the scalar "
        "the model is conditioned on, and the x-axis of every plot that follows", False),
    (0, "Bounds preserved: cov ≤ cost ≤ Σᵥ Bᵥ, so both terms stay in (0,1] before α", False),
    (0, "MILP deltas: commodity key (o,d,t0) → (o,d,t0,B); time-expanded DAG pruned at min(H, t0+B); "
        "coverage grid still on [0,H)", False),
    (0, "Verified reduction: at Bᵥ ≥ H the budget MILP returns the previous full-horizon solution "
        "link-for-link — a strict generalisation, not a different problem", True),
    (0, "Bonus: tight budgets prune the DAG, so labelling is ~30× FASTER (0.9 s vs 26 s mean per case)", False),
    (0, "Taxonomy: budget is the defining constraint of Orienteering / Team-Orienteering — this moves "
        "the benchmark into the family the routing foundation models actually transfer to", False),
], size=13.5)

# ---------------------------------------------------------------- S16
s = clone_content(prs, content_src)
set_ph(s, 11, "Extension 2")
set_ph(s, 10, "Pipeline changes")
add_bullets(s, 0.7, 1.55, 11.9, [
    (0, "① Label farm — every vehicle draws ρᵥ; fleets are heterogeneous (65%) or homogeneous (35%)", True),
    (1, "trained ρ anchors {1.0, 1.5, 2.0, 3.0} · held out for the exam: {1.25, 1.75} and {4.0}", False),
    (1, "10 780 labels, 0 errors, 0 hit the 60 s cap (previous farm: 5.1% capped) · ~17 min on 15 workers", False),
    (0, "② Encoder task token gains the budget attribute (RouteFinder-style attribute embedding):", True),
    (-1, "        Emb_o(o) + Emb_d(d) + Proj(t0)   →   … + Proj([ρᵥ , Bᵥ/H])          (+0.3 k params)", False),
    (0, "③ Decoder feasibility mask switches threshold — the one change that makes budgets GENERALISE:", True),
    (-1, "        min_finish(j, t, d) ≤ H   →   min_finish(j, t, d) ≤ t0ᵥ + Bᵥ", False),
    (1, "the budget enters through a HARD mechanical channel, not only as a soft preference in the "
        "labels ⇒ an unseen budget still yields a feasible route by construction", False),
    (0, "④ Exam gains layer L4 = budgets never seen in training (interpolation + extrapolation), plus a "
        "coverage-vs-ρ response curve against Gurobi", True),
    (0, "Unchanged: teacher forcing on next-link CE, model selection on masked-greedy objective gap, "
        "no solver in the training loop", False),
], size=13.5)

# ---------------------------------------------------------------- S17
s = clone_content(prs, content_src)
set_ph(s, 11, "Extension 2")
set_ph(s, 10, "Results — five-layer exam")
add_pic(s, f"{R}/figS3_layers.png", 1.97, 1.5, 9.4)
add_bullets(s, 0.7, 6.02, 11.9, [
    (0, f"3 seeds × 60 epochs (1.04 M params, ~4 min/seed, one RTX 3090) · 6–9 ms per case · "
        f"{AGG['L1_same']['match_mean']:.0f}/800 of L1 match or beat Gurobi", False),
    (0, "NOT comparable with the extension-1 numbers: different denominator (Σ Bᵥ vs V·H) and a "
        "different label set — a new baseline, not an improvement claim", True),
], size=12.5)

# ---------------------------------------------------------------- S18
s = clone_content(prs, content_src)
set_ph(s, 11, "Extension 2")
set_ph(s, 10, "Results — the dial generalises")
add_pic(s, f"{R}/figS2_response_curve.png", 2.17, 1.5, 9.0)
add_bullets(s, 0.7, 6.05, 11.9, [
    (0, f"Monotone response reproduced — including at ρ = 1.25 / 1.75 / 4.0, absent from every training "
        f"label. At ρ = 1 the model is EXACTLY optimal (60/60); tracking degrades as slack grows "
        f"(ρ = 4: {CURVE[-1]['model_cov']:.0f} vs {CURVE[-1]['gurobi_cov']:.0f} cells)", True),
], size=12.5)

# ---------------------------------------------------------------- closing
s = clone_content(prs, closing_src)
set_ph(s, 11, "Summary")
set_ph(s, 10, "Where this leaves us")
add_bullets(s, 0.7, 1.65, 11.9, [
    (0, "One recipe, now with a physically meaningful control: Gurobi labels offline → small "
        "Transformer → guaranteed-feasible millisecond inference, conditioned on per-vehicle budget", True),
    (0, "The α analysis is a result in its own right: a normalized weight can still be a switch when "
        "the two objective terms are marginally 1:1 coupled — worth a paragraph in the paper", False),
    (0, "Known degeneracy, deliberately left open: with uniform cell utility, overlap = 0 in 92% of "
        "labels ⇒ coverage ≈ cost. The budget controls HOW MUCH to roam; WHICH cells to prefer needs "
        "heterogeneous utility wᵢ (deferred by decision — the α₂/α₁ ratio becomes a real threshold then)", False),
    (0, "Next levers, in order: multi-sample (non-greedy) decoding · data-volume slope · wider ρ range "
        "to repair L4b extrapolation · then the 5×5 real network", False),
], size=14)

# ------------------------------------------------- reorder: new section last,
# old closing slide dropped in favour of the refreshed one
lst = prs.slides._sldIdLst
ids = list(lst)
old_closing = ids[n_v1 - 1]
lst.remove(old_closing)
prs.save(OUT)
print(f"saved {OUT} with {len(prs.slides)} slides "
      f"({n_v1 - 1} from v1 + 7 new)")
