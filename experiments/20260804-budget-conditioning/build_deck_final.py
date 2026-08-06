"""THE deck (v2, 2026-08-05) — one consolidated document, current benchmark only.

Revision per user comments (YIL-113, 2026-08-05):
  1. dedicated slide on HOW each vehicle's budget B_v is set (algorithm, not a constant)
  2. old and new objective functions listed side by side
  3. rigorous two-slide mathematical treatment of why the weight is a switch
  4. literature-review slides (setting comparison matrix + takeaway)
  5. pipeline split into TWO slides: training pipeline (with the training mechanism)
     and case-by-case generation
  6. model architecture as NATIVE PPT SHAPES (rounded boxes, elbow connectors,
     grid-aligned) across one overview + two detail slides — fully editable.

Equations are pre-rendered PNGs (render_eqs.py). Charts stay PNGs. All diagrams
are native shapes. Built from the user's template (master/theme/banner kept).

Revision 2026-08-05 later (YIL-125 r1): slide-4 caption now states the heatmap
is the (t0=0, delta=0) slice and explains the link-id numbers; two new slides
after it — TD travel-time animation (figA_td_travel.gif, animates in slideshow)
and the delta operating-day figure (figA_delta_day.png), both from build_anim.py.

Revision 2026-08-05 latest (YIL-125 r3): mod-96 cost adopted (user decision).
Env-defining pages flipped to the new law: slide 3 formula; slide 4 uses
figA_network_mod.png; slide 5 uses figA_td_travel_mod.gif (full period, step 1);
slide 6 uses figA_delta_day_mod.png (diverging: congestion can speed trips up).
Step-0 slide carries the transition note: FIFO broken at the wrap -> exact
(link, entry-time) search; labels/model/H still pre-mod until the re-run.
"""

import copy
import json
import os

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
TPL = f"{HERE}/ppt/template.pptx"
OUT = f"{HERE}/ppt/sensing_routing_method_and_results.pptx"
R = f"{HERE}/results"
EQ = f"{R}/eq"
EMU_IN = 914400

AGG = json.load(open(f"{HERE}/results_mod/agg_3seed.json"))
CURVE = json.load(open(f"{HERE}/results_mod/curve.json"))

INK, INK2 = "0B0B0B", "52514E"
BLUE, ORANGE = "2A78D6", "EB6834"
FILL, FILL_B, FILL_O = "F2F4F7", "EAF1FB", "FDF1E9"
EDGE, GREY = "C9C8C2", "8A8880"


def IN(v):
    return Emu(int(v * EMU_IN))


# ----------------------------------------------------------------- template ops
def clone(prs, src):
    new = prs.slides.add_slide(src.slide_layout)
    for shp in list(new.shapes):
        shp._element.getparent().remove(shp._element)
    for shp in src.shapes:
        new.shapes._spTree.append(copy.deepcopy(shp._element))
    return new


def set_ph(slide, idx, text):
    for sh in slide.shapes:
        if sh.is_placeholder and sh.placeholder_format.idx == idx:
            p = sh.text_frame.paragraphs[0]
            if p.runs:
                p.runs[0].text = text
                for r in p.runs[1:]:
                    r._r.getparent().remove(r._r)
            else:
                p.text = text
            return sh
    return None


def bullets(slide, x, y, w, lines, size=14):
    tb = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(1))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (lvl, txt, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " if lvl == 0 else "   – ") + txt if lvl >= 0 else txt
        p.font.size = Pt(size if lvl <= 0 else size - 1.5)
        p.font.bold = bold
    return tb


def pic(slide, path, w, y, caption=None, cap_size=12.5, x=None,
        cap_x=0.7, cap_w=11.9):
    if x is None:
        x = (13.333 - w) / 2
    slide.shapes.add_picture(path, IN(x), IN(y), width=IN(w))
    if caption:
        iw, ih = Image.open(path).size
        bullets(slide, cap_x, y + w * ih / iw + 0.10, cap_w, caption,
                size=cap_size)


def eq(slide, name, y, height=None, x=None):
    """Place an equation PNG centred at its natural 200-dpi size (or height)."""
    path = f"{EQ}/{name}.png"
    iw, ih = Image.open(path).size
    h = height if height else ih / 200
    w = h * iw / ih
    if x is None:
        x = (13.333 - w) / 2
    slide.shapes.add_picture(path, IN(x), IN(y), height=IN(h))
    return w, h


# ----------------------------------------------------------- native diagrams
def nbox(slide, x, y, w, h, title, body=None, fill=FILL, edge=EDGE, tc=INK,
         tsize=11.5, bsize=8.8, bold=True, round_=0.10):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                IN(x), IN(y), IN(w), IN(h))
    try:
        sp.adjustments[0] = round_
    except Exception:
        pass
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    sp.line.color.rgb = RGBColor.from_string(edge)
    sp.line.width = Pt(1.2)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = IN(0.06)
    tf.margin_top = tf.margin_bottom = IN(0.03)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(tsize)
    p.font.bold = bold
    p.font.color.rgb = RGBColor.from_string(tc)
    p.alignment = PP_ALIGN.CENTER
    if body:
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(bsize)
        p2.font.color.rgb = RGBColor.from_string(INK2)
        p2.alignment = PP_ALIGN.CENTER
    return sp


def _arrowhead(conn):
    ln = conn.line._get_or_add_ln()
    tail = ln.find(qn("a:tailEnd"))
    if tail is None:
        tail = ln.makeelement(qn("a:tailEnd"),
                              {"type": "triangle", "w": "med", "len": "med"})
        ln.append(tail)


def elbow(slide, a, b, fr=3, to=1, color=GREY, width=1.5, dashed=False):
    """Elbow connector glued to shapes a->b. Connection sites: 0 top, 1 left,
    2 bottom, 3 right (preset rounded rectangle)."""
    c = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, IN(1), IN(1), IN(2), IN(2))
    c.begin_connect(a, fr)
    c.end_connect(b, to)
    c.line.color.rgb = RGBColor.from_string(color)
    c.line.width = Pt(width)
    c.shadow.inherit = False
    if dashed:
        ln = c.line._get_or_add_ln()
        d = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
        ln.append(d)
    _arrowhead(c)
    return c


def label(slide, x, y, w, text, size=8.8, color=INK2, align=PP_ALIGN.CENTER,
          bold=False):
    tb = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor.from_string(color)
    p.alignment = align
    return tb


def ntable(slide, x, y, w, col_ws, rows, header_fill=FILL, font=9,
           header_font=9.5, row_h=0.32):
    """rows = list of lists of (text, bold, color) or plain str."""
    nr, nc = len(rows), len(col_ws)
    gf = slide.shapes.add_table(nr, nc, IN(x), IN(y), IN(w), IN(row_h * nr))
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = True
    for j, cw in enumerate(col_ws):
        tbl.columns[j].width = IN(cw)
    for i, row in enumerate(rows):
        for j, cell in enumerate(row):
            txt, bold, color = (cell if isinstance(cell, tuple)
                                else (cell, False, INK))
            c = tbl.cell(i, j)
            c.margin_left = c.margin_right = IN(0.05)
            c.margin_top = c.margin_bottom = IN(0.015)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]
            p.text = txt
            p.font.size = Pt(header_font if i == 0 else font)
            p.font.bold = bold or (i == 0)
            p.font.color.rgb = RGBColor.from_string(color)
            if j > 0:
                p.alignment = PP_ALIGN.CENTER
        if i == 0:
            for j in range(nc):
                tbl.cell(i, j).fill.solid()
                tbl.cell(i, j).fill.fore_color.rgb = \
                    RGBColor.from_string(header_fill)
    return tbl


def codebox(slide, x, y, w, h, lines, size=10, fill="F7F7F5"):
    """Monospace pseudocode panel (native shape, editable)."""
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                IN(x), IN(y), IN(w), IN(h))
    try:
        sp.adjustments[0] = 0.04
    except Exception:
        pass
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    sp.line.color.rgb = RGBColor.from_string(EDGE)
    sp.line.width = Pt(1.2)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = IN(0.14)
    tf.margin_top = tf.margin_bottom = IN(0.06)
    for i, (txt, hl) in enumerate(lines):
        pr = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        pr.text = txt
        pr.font.size = Pt(size)
        pr.font.name = "Courier New"
        pr.font.bold = hl
        pr.font.color.rgb = RGBColor.from_string(ORANGE if hl else INK)
        pr.alignment = PP_ALIGN.LEFT
    return sp


# ================================================================== build
prs = Presentation(TPL)
n0 = len(prs.slides)
S = list(prs.slides)
TITLE, AGENDA, CONTENT = S[0], S[1], S[3]

# ------------------------------------------------------------------ 1 title
s = clone(prs, TITLE)
set_ph(s, 0, "Yilin Wang    August 2026")
set_ph(s, 11, "Learning to Route for Sensing Coverage")
set_ph(s, 10, "A neural amortised solver for budget-constrained fleet routing "
              "— method and results")

# ------------------------------------------------------------------ 2 agenda
s = clone(prs, AGENDA)
set_ph(s, 10, "202608")
set_ph(s, 11, "Agenda")
bullets(s, 1.2, 1.9, 10.8, [
    (0, "The problem and the network", False),
    (0, "Change (a) — the objective: before vs now · why the weight is a switch "
        "(rigorous) · how each budget Bᵥ is set", True),
    (0, "Where this sits in the literature — setting comparison", False),
    (0, "Pipelines: training  ·  case-by-case generation", False),
    (0, "Change (b) — the model: overview + encoder + decoder (editable diagrams)", True),
    (0, "Evaluation design and results (3 seeds, 2 300 held-out cases)", False),
], size=16)

# ------------------------------------------------------------------ 3 problem
s = clone(prs, CONTENT)
set_ph(s, 11, "The problem")
set_ph(s, 10, "Definition")
bullets(s, 0.7, 1.6, 11.9, [
    (0, "One case (an 'operating day'): congestion δ ∈ {0,1}⁸⁰ per directed link · "
        "V vehicles · per-vehicle task (origin gate o, destination gate d, departure t₀, budget B)", False),
    (0, "Vehicles run SIMULTANEOUSLY on one timeline. Entering link i at time t occupies "
        "space–time cells (i, t … t+c−1), with travel time c(i,t) = ((base(i)+t) mod 96)//4 "
        "+ 1 + δᵢ — periodic congestion (period 96), bounded ≤ 24 + δᵢ", False),
    (0, "Objective: maximise the fleet's sensing coverage (union of occupied cells — overlap "
        "counts once), pay for travel, respect each vehicle's own budget (next slide, "
        "spelled out)", True),
    (0, "Bᵥ = the vehicle's shift length / remaining energy;  α = 0.3 / 0.7", False),
    (0, "Taxonomy: budget-constrained, selective, max-utility fleet routing = "
        "TEAM-ORIENTEERING type — not a visit-all min-cost CVRP", True),
    (0, "Gurobi solves it exactly (multi-commodity time-expanded flow) and is the source of "
        "training labels AND the benchmark; the learned model replaces per-case solving at deployment", False),
])

# ------------------------------------------------------------------ 4 network
s = clone(prs, CONTENT)
set_ph(s, 11, "The problem")
set_ph(s, 10, "Network & instances")
pic(s, f"{R}/figA_network_mod.png", 9.4, 1.5, [
    (0, "Fixed bidirectional 4×4: 80 directed links, no U-turns · 8 gates · all 56 ordered ODs "
        "feasible. Link ids 1–80 (E 1–20 · W 21–40 · S 41–60 · N 61–80) = the model's token "
        "ids; base(i) = min(id, reverse id) is each street's congestion PHASE — every street "
        "cycles the same cost range 1–24, offset by base (at t = 0: E/W cheap, N/S mid-cycle)", False),
    (0, "Heatmap = earliest arrival departing t₀ = 0 on the δ = 0 day — ONE frame of a "
        "PERIODIC quantity (period 96; next two slides)", True),
    (0, "Network and OD set FIXED; per case: δ, V, and each vehicle's (o, d, t₀, B) · horizon "
        "H = 135 (re-calibrated under mod-96, same worst-case rule; pre-mod value 338)", False),
], cap_size=11.5, cap_x=2.05, cap_w=10.6)

# ------------------------------------- 5 time-dependence (animated, YIL-125 r1)
s = clone(prs, CONTENT)
set_ph(s, 11, "The problem")
set_ph(s, 10, "Time-dependence")
pic(s, f"{R}/figA_td_travel_mod.gif", 9.8, 1.45)
bullets(s, 2.05, 5.75, 10.6, [
    (0, "The previous slide's heatmap is the t₀ = 0 frame of this loop: t₀ sweeps one FULL "
        "period 0…95 in steps of 1 and the matrix returns to its start — bounded and periodic: "
        "entry cost ≤ 24, every trip ≤ 84 steps at every t₀; G1→G2 cycles 2…43, G7→G3 "
        "cycles 34…84", True),
    (0, "Cost law (adopted 2026-08-05): c(i,t) = ((base(i)+t) mod 96)//4 + 1 + δᵢ. The wrap "
        "makes the cost non-FIFO, so these matrices come from an exact search over "
        "(link, entry-time) states — the Step 0 slide says what this changes in the machinery", False),
    (0, "Animated GIF — plays in slideshow mode; in print this page shows the t₀ = 0 frame", False),
], size=11.5)

# ------------------------------------------- 6 day-to-day delta (YIL-125 r1)
s = clone(prs, CONTENT)
set_ph(s, 11, "The problem")
set_ph(s, 10, "Day-to-day: δ")
pic(s, f"{R}/figA_delta_day_mod.png", 9.6, 1.5)
bullets(s, 2.05, 5.55, 10.6, [
    (0, "δ ∈ {0,1}⁸⁰: one bit per DIRECTED link, drawn i.i.d. Bernoulli(½) per operating day "
        "(the farm's sampler); δᵢ = 1 adds +1 to EVERY traversal of link i that day", True),
    (0, "An instance property shared by the whole fleet — not a vehicle attribute; the two "
        "directions of one street can differ (left: one-way congestion is visible)", False),
    (0, "Effect at t₀ = 0 (right): this sampled day shifts the 56 ODs by −9 … +7 steps — under "
        "mod-96, congestion can even make a trip FASTER (4/56 here): a +1 delay can push a "
        "later link past its wrap into the cheap zone. The model sees δ per link token: "
        "link_emb(i) + Proj(δᵢ) (encoder slide)", False),
], size=11.5)

# ------------------------------------------ 5 change (a): objective before/after
s = clone(prs, CONTENT)
set_ph(s, 11, "Change (a) — objective")
set_ph(s, 10, "Objective, before vs now")
label(s, 0.75, 1.55, 6.0, "BEFORE — weight-controlled, one global horizon H",
      size=12, color=INK2, align=PP_ALIGN.LEFT, bold=True)
eq(s, "eq_old", 1.95, height=0.62)
label(s, 0.75, 2.85, 8.0, "NOW — budget-controlled, per-vehicle deadline "
      "(objective terms unchanged)", size=12, color=ORANGE,
      align=PP_ALIGN.LEFT, bold=True)
eq(s, "eq_new", 3.25, height=0.72)
ntable(s, 1.35, 4.45, 10.6, [2.6, 4.0, 4.0], [
    [("what changed", True, INK), ("before", False, INK), ("now", True, INK)],
    ["control variable", "weight α₂ (a preference)",
     ("per-vehicle budget Bᵥ (a physical resource)", True, INK)],
    ["hard constraint", "arrive by the global horizon H",
     ("arrive by own deadline t₀ᵥ + Bᵥ", True, INK)],
    ["normaliser", "V · H  (fleet × horizon)",
     ("Σᵥ Bᵥ  (available driving time); cov ≤ cost ≤ Σᵥ Bᵥ keeps both terms in (0,1]", False, INK)],
    ["solution response", "step function of α₂ (next 3 slides)",
     ("monotone in Bᵥ — a usable dial", True, INK)],
], row_h=0.42)
bullets(s, 2.05, 6.72, 10.7, [
    (0, "Verified reduction: at Bᵥ ≥ H the new MILP returns the old one's solution "
        "link-for-link — a strict generalisation, not a different problem", False),
], size=11.5)

# ------------------------------------------------ 6 change (a): rigorous math I
s = clone(prs, CONTENT)
set_ph(s, 11, "Change (a) — objective")
set_ph(s, 10, "Why α is a switch — proof I")
label(s, 0.75, 1.50, 11.8, "Setup.  For a fixed instance (network, δ, tasks, "
      "deadlines), a feasible solution R = (R₁ … R_V):", size=12.5,
      align=PP_ALIGN.LEFT)
eq(s, "eq_def", 1.90, height=0.52)
label(s, 0.75, 2.62, 11.8, "Lemma (cell accounting).  Each time-step of travel "
      "occupies exactly one (link, time) cell, and the union counts every cell "
      "once. Hence", size=12.5, align=PP_ALIGN.LEFT)
eq(s, "eq_lemma", 3.06, height=0.34)
label(s, 0.75, 3.55, 11.8, "Measured on all 10 780 optimal labels: median K/C = "
      "0.999; K = C exactly (zero overlap) in 84–93 % per shard — the "
      "zero-overlap regime is the TYPICAL one, not a corner case.", size=11,
      color=INK2, align=PP_ALIGN.LEFT)
label(s, 0.75, 4.30, 11.8, "Theorem (collapse in the zero-overlap regime).  "
      "If K(R) = C(R) for the instance's candidate solutions, the objective "
      "degenerates to a single scalar multiple of total travel:", size=12.5,
      align=PP_ALIGN.LEFT, bold=True)
eq(s, "eq_collapse", 4.88, height=0.52)
bullets(s, 0.7, 5.65, 12.0, [
    (0, "Minimising it therefore means:  α₂ < α₁ → minimise C (min-time routes) ·  "
        "α₂ > α₁ → MAXIMISE C (fill the budget) ·  α₂ = α₁ → every feasible "
        "solution ties (total degeneracy)", True),
    (0, "The map  α₂ ↦ optimal-solution set  is a STEP function with its only jump at "
        "α₂/α₁ = 1. No value of α₂ selects anything between the two extremes — "
        "and this is independent of any normaliser N", True),
    (0, "α's only residual role: for α₂ < α₁ it breaks ties AMONG min-cost solutions "
        "(more coverage preferred at equal cost)", False),
], size=12)

# ----------------------------------------------- 7 change (a): rigorous math II
s = clone(prs, CONTENT)
set_ph(s, 11, "Change (a) — objective")
set_ph(s, 10, "Why α is a switch — proof II")
label(s, 0.75, 1.50, 11.8, "Proposition A (the cost regime is robust).  Let "
      "C* = min feasible cost and K₀ = best coverage among min-cost solutions. "
      "Using K ≤ C:", size=12.5, align=PP_ALIGN.LEFT)
eq(s, "eq_prop1", 1.94, height=0.52)
bullets(s, 0.7, 2.62, 12.0, [
    (0, "Proof: f(R) ≥ (α₁−α₂)C(R)/N for every R, while the best min-cost solution "
        "attains (α₁C* − α₂K₀)/N; compare. If a min-cost solution has zero overlap "
        "(K₀ = C*), the threshold is 0: EVERY optimum is min-cost — overlap only "
        "shifts the switch point, it cannot create interior positions", False),
], size=11.5)
label(s, 0.75, 3.42, 11.8, "Proposition B (the roam regime saturates the "
      "budget).  A deadline-feasible detour of length k with ΔK fresh cells "
      "changes the objective by", size=12.5, align=PP_ALIGN.LEFT)
eq(s, "eq_prop2", 3.86, height=0.50)
bullets(s, 0.7, 4.52, 12.0, [
    (0, "For α₂ > α₁ every FULLY-fresh detour (ΔK = k) strictly improves — so at any "
        "optimum, no vehicle can still afford one: budgets bind. The threshold "
        "ΔK/k > α₁/α₂ also explains why solutions barely move across α₂ ∈ "
        "[0.52, 0.90] in the zero-overlap regime", False),
    (0, "Conclusion: α selects one of two regimes; the quantity that positions a "
        "solution INSIDE the roam regime is the budget. C*(B) grows monotonically in "
        "bounded increments as B grows — a graded, learnable response", True),
    (0, "That is change (a): keep α fixed (0.3/0.7, roam regime), move the control to "
        "the constraint — the per-vehicle budget Bᵥ. Next slide: both claims measured "
        "on one instance", False),
], size=12)

# ------------------------------------------------ 8 change (a): measured
s = clone(prs, CONTENT)
set_ph(s, 11, "Change (a) — objective")
set_ph(s, 10, "Switch vs dial — measured")
pic(s, f"{R}/figS1_knob_vs_switch.png", 11.6, 1.5, [
    (0, "One fixed instance, one MILP, one solver, MIPGap 0.5 %. LEFT: 13 values of α₂ with the budget "
        "fixed at H → 2 distinct solutions (α₂ = 0.10–0.48 identical link-for-link). "
        "RIGHT: α frozen at 0.3/0.7, per-vehicle budget swept → 10 distinct, monotone solutions", False),
])

# ------------------------------------------------ 9 TD-Dijkstra (prerequisite)
s = clone(prs, CONTENT)
set_ph(s, 11, "Change (a) — objective")
set_ph(s, 10, "Step 0: earliest arrival")
label(s, 0.75, 1.48, 11.9, "Before any budget can be set, we must know the fastest "
      "this vehicle could possibly make its trip TODAY — earliest arrival on the "
      "current network, under this case's congestion δ. One run per vehicle:",
      size=12.5, align=PP_ALIGN.LEFT, bold=True)
codebox(s, 0.75, 2.14, 9.10, 2.75, [
    ("EarliestArrival(o, t₀, δ):   # one run per vehicle · exact under mod-96", True),
    ("  push (t₀, l)  for every link l leaving gate o", False),
    ("  while queue not empty:", False),
    ("      (t, l) ← pop the SMALLEST entry time    # state = (link, entry time)", False),
    ("      if (l, t) already expanded: continue    # dedup STATES, not links", True),
    ("      s ← t + c(l, t)      # exit time; the clock is INSIDE the cost", True),
    ("      if l ends at a gate g:  arrive[g] ← min(arrive[g], s)", False),
    ("      for every successor j of l (no U-turn):  push (s, j)", False),
    ("  return arrive                               # earliest arrival at EVERY gate", False),
    ("", False),
    ("τᵐⁱⁿ_v = arrive[d_v] − t₀_v", True),
], size=9.6)
b1 = nbox(s, 10.10, 2.14, 2.70, 2.75, "Where it runs",
          "① budget assignment:\nV runs per case\n(~41 000 in the farm)\n\n"
          "② decoder mask:\nmin_finish(l, t, d) —\nsame algorithm from one\nlink, memoised\n\n"
          "③ horizon calibration:\nonce, all 8 gates",
          fill=FILL_O, edge=ORANGE, tc=ORANGE, tsize=10.5, bsize=8.6)
bullets(s, 2.05, 5.15, 10.55, [
    (0, "Ordinary Dijkstra with two changes: a label is a TIME (today's congestion δ sits "
        "inside c), and there is one label per (link, entry-time) STATE, popped in time order — "
        "exact for any positive time-dependent cost, FIFO or not", True),
    (0, "Why states are necessary: the mod-96 wrap breaks FIFO — entering a link one step later "
        "can mean exiting 22 steps earlier — so 'the first label per LINK is final' fails "
        "(single-label TD-Dijkstra is wrong on 18 % of (OD, t₀) queries under this cost); "
        "per-state dedup restores the classic invariant", True),
    (0, "One run = bounded states on 80 links → sub-millisecond; ~41 000 runs across the farm "
        "are negligible next to the MILP solves, and the mask's memoised state-searches sit "
        "inside the per-case inference time. All numbers in this deck are the mod-96 benchmark "
        "(H = 135, re-run 2026-08-06)", False),
], size=11.5)

# ------------------------------------------------ 10 how B_v is set
s = clone(prs, CONTENT)
set_ph(s, 11, "Change (a) — objective")
set_ph(s, 10, "How Bᵥ is set")
label(s, 0.75, 1.48, 11.9, "In plain words: first compute the fastest this vehicle "
      "could make its trip today (τᵐⁱⁿᵥ, previous slide); the budget is that minimum "
      "times a slack factor — 'you get ρ× the minimum time'. Not a fixed constant: "
      "derived per vehicle, per instance.", size=12.5,
      align=PP_ALIGN.LEFT, bold=True)
eq(s, "eq_budget", 2.06, height=0.40)
eq(s, "eq_taumin", 2.58, height=0.36)
bullets(s, 0.7, 3.12, 12.0, [
    (0, "Step 1 — draw the slack ratio ρᵥ.  Training farm: per vehicle from the anchors "
        "{1.0, 1.5, 2.0, 3.0}; in 35 % of fleets all vehicles share one ρ (homogeneous), "
        "in 65 % each vehicle draws its own (heterogeneous). At TEST time ρ is arbitrary — "
        "including values never seen in training ({1.25, 1.75} interpolation, 4.0 extrapolation)", False),
    (0, "Step 2 — the feasibility floor τᵐⁱⁿᵥ comes from the TD-Dijkstra run of the previous "
        "slide (earliest arrival o→d departing at t₀ under THIS instance's δ) — "
        "instance-specific, not a table lookup", False),
    (0, "Step 3 — Bᵥ = ⌈ρᵥ · τᵐⁱⁿᵥ⌉, floored at τᵐⁱⁿᵥ (a task is never born infeasible) and "
        "capped at H − t₀ᵥ (never beyond the coverage grid)", False),
    (0, "Why a RATIO and not an absolute number: the same B = 50 is vacuous for a near OD "
        "(τᵐⁱⁿ = 14) and infeasible for a far one (τᵐⁱⁿ = 61). ρ is comparable across all "
        "ODs:  ρ = 1 ⇔ exactly the min-time trip;  ρ = 2 ⇔ twice the minimum time", True),
    (0, "Worked example (the sweep instance, V = 3, τᵐⁱⁿ = [38, 61, 14]):  ρ = 1.5 for all "
        "⇒  B = [57, 92, 21] — same slack semantics, very different absolute budgets", False),
    (0, "Bᵥ then enters BOTH sides of the pipeline: the MILP constraint (labels) and the "
        "model input (conditioning) + the decoder feasibility mask", True),
], size=12.5)

# ------------------------------------------------ 10 literature matrix
s = clone(prs, CONTENT)
set_ph(s, 11, "Literature")
set_ph(s, 10, "Setting comparison")
CK = ("✓", True, "1BAF7A")
NO = ("—", False, GREY)
PART = ("part", False, INK2)
ntable(s, 0.55, 1.55, 12.25,
       [3.05, 1.30, 1.42, 1.42, 1.28, 1.42, 1.28, 1.08],
       [
        [("work", True, INK), "selective\n(not visit-all)", "per-vehicle\n(o,d,t₀,B)",
         "arc-level fixed\nnetwork", "time-dep.\ntravel", "space–time\ncoverage",
         "learned\namortised", "budget\n0-shot"],
        ["Kool AM '19 (OP)", CK, NO, NO, NO, NO, CK, NO],
        ["TOP-Former '25 (TOP)", CK, ("budget only,\nhomog.", False, INK2), NO, NO, NO, CK, NO],
        ["UAS-MSTOP '23", CK, ("start + fuel", False, INK2), NO, NO, NO, CK, NO],
        ["MTPOMO / MVMoE / RouteFinder", NO, ("L / O attrs", False, INK2), NO, NO, NO, CK,
         ("variant-\nlevel", False, INK2)],
        ["GOAL '25 (incl. OP task)", PART, NO, NO, NO, NO, CK, NO],
        ["FM-MCVRP '24 (fixed graph SL)", NO, NO, ("fixed graph,\nnodes", False, INK2), NO, NO, CK, NO],
        ["SED2AM '25 (TD-VRP)", NO, NO, ("road data", False, INK2), CK, NO, CK, NO],
        ["Neural CARP '19–'26", NO, NO, CK, NO, NO, CK, NO],
        ["Drive-by sensing OR (Han'24, Chen'24, Zhu'14)", CK, PART, CK, PART, CK, NO, NO],
        [("OURS (benchmark v1)", True, ORANGE), CK, CK, CK, CK, CK, CK, CK],
       ], font=8.6, header_font=8.8, row_h=0.44)
bullets(s, 2.05, 6.62, 10.7, [
    (0, "No existing work is identical: every neighbour misses ≥ 2 columns. The application shelf "
        "(bottom) has the right problem but no learning; the neural shelves have learning but not "
        "this problem", True),
], size=11)

# ------------------------------------------------ 11 literature takeaway
s = clone(prs, CONTENT)
set_ph(s, 11, "Literature")
set_ph(s, 10, "Neighbours vs ours")
bullets(s, 0.7, 1.55, 12.0, [
    (0, "Closest three, and exactly what they miss:", True),
    (1, "TOP-Former (closest neural): team + per-vehicle budget MASKING — but one shared depot, "
        "homogeneous budget, synchronized start, Euclidean nodes, static prizes, RL", False),
    (1, "FM-MCVRP (closest recipe): fixed graph + SL + joint fleet sequence — but visit-all "
        "min-cost CVRP, no budgets, no time dependence, deliberately sub-optimal labels", False),
    (1, "Han et al. '24 TR-B (closest problem): space–time sensing utility + budget on real "
        "networks — but pure OR, per-instance solving, no learned router at all", False),
    (0, "Honesty box — every ingredient of our network exists somewhere: mask kernel (Kool), "
        "attribute conditioning (MTPOMO / RouteFinder), fixed-graph SL + joint sequence "
        "(FM-MCVRP), centralized team decoding (TOP-Former)", True),
    (0, "What is ours: (C1) the first learned amortised solver for THIS problem class — "
        "budget-constrained team sensing-coverage routing, arc-level, time-dependent — the "
        "bridge between the two shelves", True),
    (1, "(C2) the switch-vs-dial analysis: exact characterisation of when a scalarisation "
        "weight cannot position solutions, and the constructive replacement (budget in the "
        "constraint)", False),
    (1, "(C3) value-level zero-shot evaluation: unseen CONTINUOUS budgets (interp + extrap) + "
        "a solver-tracked response curve — extends the variant-level zero-shot of "
        "MTPOMO / RouteFinder", False),
    (0, "Threats we track: uniform-utility degeneracy (wᵢ pending) · toy scale (80 links) · "
        "greedy-only decoding · fast-moving neighbours (DeCoST ICLR'26)", False),
], size=12.5)

# ------------------------------------------------ 12 training pipeline (native)
s = clone(prs, CONTENT)
set_ph(s, 11, "Pipeline 1 / 2")
set_ph(s, 10, "Training (offline, run once)")
w5, gap = 2.24, 0.20
xs = [0.55 + i * (w5 + gap) for i in range(5)]
yA, hA = 1.55, 1.10
bA = [
    nbox(s, xs[0], yA, w5, hA, "Environment",
         "fixed 4×4 · 80 links\nc(i,t) time-dependent\nTD-Dijkstra oracle"),
    nbox(s, xs[1], yA, w5, hA, "Case sampler",
         "δ ~ {0,1}⁸⁰ · V ∈ {2,3,4,6}\nODs stratified + 4 held out\nt₀ ∈ {0…5}"),
    nbox(s, xs[2], yA, w5, hA, "Budget assignment",
         "ρᵥ ~ anchors {1, 1.5, 2, 3}\nBᵥ = ⌈ρᵥ·τᵐⁱⁿᵥ⌉\n65 % heterogeneous fleets",
         fill=FILL_O, edge=ORANGE),
    nbox(s, xs[3], yA, w5, hA, "Gurobi MILP",
         "time-expanded flow\nper-vehicle deadline t₀+B\nMIPGap 2 % · 0.90 s/case",
         fill=FILL_B, edge=BLUE, tc=BLUE),
    nbox(s, xs[4], yA, w5, hA, "Label set",
         "10 780 (case, routes)\n0 errors · 0 hit the cap\n~17 min · 15 workers"),
]
for a, b in zip(bA, bA[1:]):
    elbow(s, a, b, 3, 1)
w4, gap4 = 2.86, 0.22
x4 = [0.55 + i * (w4 + gap4) for i in range(4)]
yB, hB = 3.35, 1.10
bB = [
    nbox(s, x4[0], yB, w4, hB, "Tokenise",
         "[BOS v₁-links SEP v₂-links\nSEP … EOS] · vocab 84\none string = whole fleet"),
    nbox(s, x4[1], yB, w4, hB, "Teacher forcing",
         "predict the next link token\ncross-entropy (PAD ignored)\nAdamW 3e-4 · batch 32 · 60 ep"),
    nbox(s, x4[2], yB, w4, hB, "Select by GENERATION",
         "each epoch: masked-greedy\ndecode 150 val cases →\nTRUE objective gap",
         fill=FILL_O, edge=ORANGE),
    nbox(s, x4[3], yB, w4, hB, "Best checkpoint",
         "1.04 M params\n3 seeds · ~4 min/seed\n(one RTX 3090)"),
]
for a, b in zip(bB, bB[1:]):
    elbow(s, a, b, 3, 1)
elbow(s, bA[4], bB[0], 2, 0)
bullets(s, 0.7, 4.85, 12.0, [
    (0, "Loss is imitation only: next-link cross-entropy on Gurobi's sequences — no solver, "
        "no decoding, no reward inside the training loop", False),
    (0, "Model selection is NOT the training loss: CE keeps improving while true solution "
        "quality can worsen (imitation overfitting) — selecting on generated-solution gap "
        "catches this", True),
    (0, "The budget enters twice already at training time: as an input feature ([ρᵥ, Bᵥ/H] in "
        "the task token) and inside the validation decoder's feasibility mask", False),
    (0, "Everything above runs once; retraining is ~20 min end-to-end (farm + 3 seeds)", False),
], size=12)

# --------------------------------------- 13 case-by-case generation (native)
s = clone(prs, CONTENT)
set_ph(s, 11, "Pipeline 2 / 2")
set_ph(s, 10, "Generation (online, per case)")
yC, hC = 1.75, 1.20
c1 = nbox(s, 0.55, yC, 2.35, hC, "New case",
          "today's δ\nV tasks (o, d, t₀, B)")
c2 = nbox(s, 3.25, yC, 2.75, hC, "Encoder — runs ONCE",
          "80 link tokens + V task\ntokens → case memory")
c3 = nbox(s, 6.35, yC, 3.55, hC, "Autoregressive step",
          "decoder → logits over 84\n+ feasibility mask (−∞)\n→ argmax next link",
          fill=FILL_O, edge=ORANGE)
c4 = nbox(s, 10.55, yC, 2.25, hC, "V routes",
          "feasible by\nconstruction")
elbow(s, c1, c2, 3, 1)
elbow(s, c2, c3, 3, 1)
elbow(s, c3, c4, 3, 1)
loop = nbox(s, 6.65, 3.45, 3.55, 0.62, "⟲  append token · t += c(l, t) · repeat until EOS",
            None, fill="F7F7F5", edge=GREY, tsize=9.5, bold=False)
elbow(s, c3, loop, 2, 0, color=GREY)
mask = nbox(s, 1.0, 3.3, 4.5, 1.55, "Feasibility mask (0 parameters)",
            "① not a road successor / U-turn\n"
            "② cannot reach OWN destination by t₀ᵥ + Bᵥ  (TD-Dijkstra query)\n"
            "③ SEP only when parked next to own gate; EOS after V SEPs",
            fill=FILL_O, edge=ORANGE, tc=ORANGE, bsize=9.2)
elbow(s, mask, c3, 0, 2, color=ORANGE)
bullets(s, 0.7, 5.30, 12.0, [
    (0, "6–9 ms per case on one GPU — including the mask's earliest-arrival queries; no solver "
        "call anywhere at deployment", True),
    (0, "The mask is why unseen budgets stay feasible: feasibility is ENFORCED by rules, the "
        "network only ranks the legal moves (the budget has a hard mechanical channel, not just "
        "a soft preference from the labels)", True),
    (0, "SEP ends a vehicle; decoding continues with the next vehicle from its own (o, t₀) — "
        "later vehicles see earlier routes through self-attention (coordination)", False),
    (0, "Today: greedy decoding. Multi-sample (nucleus) decoding is the known next lever — "
        "inference-only, no retraining", False),
], size=12)

# ------------------------------------------------ 14 model overview (native)
s = clone(prs, CONTENT)
set_ph(s, 11, "Change (b) — model")
set_ph(s, 10, "Architecture overview")
m1 = nbox(s, 0.55, 1.80, 2.45, 1.50, "Inputs",
          "congestion δ (80)\nV tasks\n(o, d, t₀, ρ, B/H)")
m2 = nbox(s, 3.45, 1.80, 2.55, 1.50, "Transformer encoder",
          "×3 layers · d = 128\n4 heads · pre-norm\nruns once per case")
m3 = nbox(s, 6.45, 1.80, 2.45, 1.50, "Case memory",
          "(80 + V) × d", fill=FILL_B, edge=BLUE, tc=BLUE)
m4 = nbox(s, 9.55, 1.80, 3.20, 1.50, "Autoregressive decoder",
          "joint fleet sequence\n×3 blocks → head (84)\n→ masked softmax")
m5 = nbox(s, 9.55, 4.10, 3.20, 1.30, "Feasibility mask",
          "topology · OWN deadline\nt₀ᵥ + Bᵥ · SEP/EOS grammar",
          fill=FILL_O, edge=ORANGE, tc=ORANGE)
m6 = nbox(s, 5.60, 4.10, 2.90, 1.30, "Fleet plan",
          "V feasible routes\n6–9 ms per case")
elbow(s, m1, m2, 3, 1)
elbow(s, m2, m3, 3, 1)
elbow(s, m3, m4, 3, 1)
elbow(s, m5, m4, 0, 2, color=ORANGE)
elbow(s, m4, m6, 2, 3)
label(s, 6.30, 3.42, 2.75, "the decoder cross-attends to the\ncase memory at EVERY step", size=8.6, color=INK2)
bullets(s, 0.7, 5.85, 12.0, [
    (0, "The two redesign points (details on the next two slides):", True),
    (1, "① instance attributes (o, d, t₀, B) moved OUT of the weights and INTO the input "
        "— one model serves every OD, fleet size, departure and budget", False),
    (1, "② the mask threshold is each vehicle's OWN deadline t₀ᵥ + Bᵥ — unseen budgets "
        "stay feasible by construction", False),
    (0, "All diagram elements are native shapes — fully editable (boxes, elbow connectors)", False),
], size=12)

# ------------------------------------------------ 15 encoder detail (native)
s = clone(prs, CONTENT)
set_ph(s, 11, "Change (b) — model")
set_ph(s, 10, "Encoder — token construction")
label(s, 0.7, 1.50, 4.6, "Link tokens  (×80, one per directed link)", size=11.5,
      bold=True, align=PP_ALIGN.LEFT)
l1 = nbox(s, 0.70, 1.85, 1.95, 0.62, "Emb(link id)", None, tsize=10, bold=False)
l2 = nbox(s, 2.95, 1.85, 1.95, 0.62, "W · δᵢ", None, tsize=10, bold=False)
lt = nbox(s, 1.55, 2.90, 2.75, 0.62, "link token  (d)", None, tsize=10.5)
elbow(s, l1, lt, 2, 0)
elbow(s, l2, lt, 2, 0)
label(s, 5.6, 1.50, 7.2, "Task tokens  (×V, one per vehicle — "
      "redesign ①: attributes as INPUT)", size=11.5, bold=True,
      align=PP_ALIGN.LEFT, color=ORANGE)
t1 = nbox(s, 5.60, 1.85, 1.62, 0.62, "Emb_o(o)", None, tsize=10, bold=False)
t2 = nbox(s, 7.42, 1.85, 1.62, 0.62, "Emb_d(d)", None, tsize=10, bold=False)
t3 = nbox(s, 9.24, 1.85, 1.62, 0.62, "Proj(t₀)", None, tsize=10, bold=False)
t4 = nbox(s, 11.06, 1.85, 1.72, 0.62, "Proj([ρ, B/H])", None, tsize=10,
          bold=False, fill=FILL_O, edge=ORANGE, tc=ORANGE)
tt = nbox(s, 7.85, 2.90, 2.75, 0.62, "task token  (d)", None, tsize=10.5,
          fill=FILL_O, edge=ORANGE)
for b_ in (t1, t2, t3, t4):
    elbow(s, b_, tt, 2, 0)
label(s, 1.55, 2.62, 2.75, "sum", size=8.5)
label(s, 7.85, 2.62, 2.75, "sum (4 attribute channels)", size=8.5)
cc = nbox(s, 3.65, 4.00, 6.0, 0.60, "concatenate:  (80 + V) tokens",
          None, tsize=10.5, bold=False)
elbow(s, lt, cc, 2, 0)
elbow(s, tt, cc, 2, 0)
enc = nbox(s, 3.65, 5.00, 6.0, 0.72, "Transformer encoder  ×3",
           "self-attention + FFN · d = 128 · 4 heads · pre-norm", bsize=8.8)
elbow(s, cc, enc, 2, 0)
mem = nbox(s, 3.65, 6.10, 6.0, 0.60, "case memory   (80 + V) × d", None,
           fill=FILL_B, edge=BLUE, tc=BLUE, tsize=10.5)
elbow(s, enc, mem, 2, 0)
bullets(s, 2.05, 6.80, 10.7, [
    (0, "Why ①: one set of weights serves every OD pair, fleet size, departure and budget; the "
        "encoder reasons about the fleet JOINTLY. Same recipe as the routing foundation models "
        "(RouteFinder / MTPOMO attribute embeddings) — our budget plays their capacity's role", False),
], size=10.5)

# ------------------------------------------------ 16 decoder detail (native)
s = clone(prs, CONTENT)
set_ph(s, 11, "Change (b) — model")
set_ph(s, 10, "Decoder — one step")
toks = ["BOS", "l₁", "l₂", "SEP", "l₁", "…"]
tw, tx0, ty = 0.78, 3.30, 1.55
tok_shapes = []
for i, t in enumerate(toks):
    f, e, c = (FILL_O, ORANGE, ORANGE) if t == "SEP" else (FILL, EDGE, INK)
    tok_shapes.append(nbox(s, tx0 + i * (tw + 0.10), ty, tw, 0.52, t, None,
                           fill=f, edge=e, tc=c, tsize=9.5, bold=False))
label(s, tx0, 1.22, 7.0, "partial JOINT fleet sequence — SEP parks a vehicle, "
      "the next one starts from its own (o, t₀)", size=9,
      align=PP_ALIGN.LEFT)
d1 = nbox(s, 3.60, 2.45, 4.60, 0.58, "masked self-attention (causal)", None,
          tsize=10, bold=False)
d2 = nbox(s, 3.60, 3.35, 4.60, 0.58, "cross-attention → case memory", None,
          tsize=10, bold=False)
d3 = nbox(s, 3.60, 4.25, 4.60, 0.58, "FFN     (decoder block ×3)", None,
          tsize=10, bold=False)
d4 = nbox(s, 3.60, 5.15, 4.60, 0.58, "linear head → logits over vocab (84)",
          None, tsize=10, bold=False)
d5 = nbox(s, 3.60, 6.05, 4.60, 0.58, "softmax → next link (argmax / sample)",
          None, tsize=10, bold=False)
elbow(s, tok_shapes[3], d1, 2, 0)
elbow(s, d1, d2, 2, 0)
elbow(s, d2, d3, 2, 0)
elbow(s, d3, d4, 2, 0)
elbow(s, d4, d5, 2, 0)
memstub = nbox(s, 0.70, 3.35, 2.20, 0.58, "case memory", None,
               fill=FILL_B, edge=BLUE, tc=BLUE, tsize=10, bold=False)
elbow(s, memstub, d2, 3, 1, color=BLUE)
mk = nbox(s, 9.05, 3.90, 3.70, 1.90, "feasibility mask — redesign ②",
          "logits[illegal] = −∞\n① not a successor / U-turn\n"
          "② cannot reach own dest by t₀ᵥ + Bᵥ\n③ SEP only next to own gate",
          fill=FILL_O, edge=ORANGE, tc=ORANGE, bsize=9.2)
elbow(s, mk, d4, 1, 3, color=ORANGE)
elbow(s, d5, tok_shapes[5], 3, 2, color=GREY)
label(s, 8.42, 2.42, 4.2, "append next token · advance clock t += c(l, t) · repeat until EOS",
      size=8.4, color=INK2, align=PP_ALIGN.LEFT)
bullets(s, 2.05, 6.80, 10.7, [
    (0, "Why ②: the budget gets a HARD channel into decoding (mask), not only a soft preference "
        "from labels — feasibility is enforced, not learned, so budgets never seen in training "
        "still decode to feasible routes. The mask has no parameters", False),
], size=10.5)

# ------------------------------------------------ 17 label set
s = clone(prs, CONTENT)
set_ph(s, 11, "Data")
set_ph(s, 10, "Label set")
bullets(s, 0.7, 1.65, 11.9, [
    (0, "10 780 Gurobi labels · 0 errors · 0 solves hit the 60 s cap · 0.90 s mean per case "
        "(tight budgets prune the time-expanded graph → labelling ~30× faster than the "
        "full-horizon farm) · ~17 min on 15 workers, fully resumable", True),
    (0, "Shards:", True),
    (1, "train 8 000 · same-distribution test 800", False),
    (1, "OD zero-shot 400 (4 gate pairs never trained) · fleet extrapolation 400 (V ∈ {5,8})", False),
    (1, "UNSEEN-budget interpolation 400 (ρ ∈ {1.25, 1.75}) · extrapolation 300 (ρ = 4.0)", False),
    (1, "response curve: 60 fixed instances × 8 ρ values, re-solved by Gurobi at every ρ", False),
    (0, "Trained ρ anchors {1.0, 1.5, 2.0, 3.0}; the three held-out ρ values appear in NO "
        "training label", False),
    (0, "Every label passes two checks at write time: flow-decomposed objective ≡ Gurobi ObjVal "
        "(1e-6) and every per-vehicle deadline re-validated", False),
], size=13.5)

# ------------------------------------------------ 18 eval design
s = clone(prs, CONTENT)
set_ph(s, 11, "Evaluation")
set_ph(s, 10, "Design")
bullets(s, 0.7, 1.65, 11.9, [
    (0, "Five layers, 2 300 cases, none seen in training:", True),
    (1, "L1 same-distribution — 800 fresh cases (trained ODs, new δ / task mixes)", False),
    (1, "L2 OD ZERO-SHOT — 400 cases on 4 OD pairs never trained (G1→G5, G6→G2, G4→G8, G7→G3)", False),
    (1, "L3 fleet extrapolation — 400 cases with V ∈ {5, 8}, trained on {2,3,4,6}", False),
    (1, "L4a UNSEEN budget, interpolation — 400 cases at ρ ∈ {1.25, 1.75}", False),
    (1, "L4b UNSEEN budget, extrapolation — 300 cases at ρ = 4.0, beyond the trained range", False),
    (0, "Plus the RESPONSE CURVE: 60 fixed instances re-solved at every ρ by both Gurobi and "
        "the model — did it learn the budget as a continuous quantity or memorise four settings?", True),
    (0, "Metrics: feasibility rate · objective gap vs Gurobi (absolute + relative) · fraction "
        "matching or beating Gurobi · inference time · per-case CSV. 3 seeds, mean ± std", False),
    (0, "Tuning used a validation slice of the training shard only; every reported number is "
        "from held-out shards", False),
], size=13.5)

# ------------------------------------------------ 19 results layers
s = clone(prs, CONTENT)
set_ph(s, 11, "Results")
set_ph(s, 10, "Five-layer exam")
pic(s, f"{R}/figS3_layers.png", 9.4, 1.5, [
    (0, f"3 seeds × 60 epochs · 2 300 held-out cases, 100 % feasible in every layer · 6–9 ms per case "
        f"· {AGG['L1_same']['match_mean']:.0f}/800 of L1 match or beat Gurobi", False),
    (0, "Unseen budgets (L4a, orange) cost about as much accuracy as an unseen fleet size — the "
        "budget generalises like the other conditioned attributes, not worse", True),
])

# ------------------------------------------------ 20 results curve
s = clone(prs, CONTENT)
set_ph(s, 11, "Results")
set_ph(s, 10, "Budget response curve")
pic(s, f"{R}/figS2_response_curve.png", 9.0, 1.5, [
    (0, f"Monotone response reproduced — including at ρ = 1.25 / 1.75 / 4.0, absent from every training "
        f"label. At ρ = 1 the model is EXACTLY optimal (60/60); tracking degrades as slack grows "
        f"(ρ = 4: {CURVE[-1]['model_cov']:.0f} vs {CURVE[-1]['gurobi_cov']:.0f} cells) — the limit "
        f"of greedy decoding", True),
])

# ------------------------------------------------ 21 summary
s = clone(prs, CONTENT)
set_ph(s, 11, "Summary")
set_ph(s, 10, "Where this stands")
bullets(s, 0.7, 1.7, 11.9, [
    (0, "One recipe, verified end to end: Gurobi labels offline → 1.04 M-parameter Transformer → "
        "guaranteed-feasible millisecond inference, conditioned on each vehicle's own budget", True),
    (0, "The control variable is now physical (shift length / energy), continuous, and demonstrably "
        "generalises to budgets never trained on", False),
    (0, "Positioning: the first learned amortised solver for this problem class — the bridge between "
        "the drive-by-sensing OR shelf (right problem, no learning) and the neural routing shelf "
        "(right method, different problem)", False),
    (0, "Verification chain at every step: MILP ≡ flow decomposition ≡ simulator · budget MILP ≡ the "
        "full-horizon MILP at B ≥ H · per-case CSVs against Gurobi on every held-out layer", False),
    (0, "Known and deliberately open: with uniform cell utility, overlap = 0 in ~92 % of optimal "
        "solutions, so coverage ≈ cost. The budget controls HOW MUCH to roam; making WHICH cells a "
        "real decision needs heterogeneous utility wᵢ — the next modelling step, not a fix", False),
    (0, "Next levers, in order: multi-sample (non-greedy) decoding · wider ρ range for extrapolation "
        "· data-volume slope · heterogeneous wᵢ · then the 5×5 real network", True),
], size=13.5)

xml = prs.slides._sldIdLst
for sld in list(xml)[:n0]:
    xml.remove(sld)
prs.save(OUT)
print(f"saved {OUT} with {len(prs.slides)} slides")
