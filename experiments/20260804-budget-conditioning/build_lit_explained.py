"""Standalone literature-review explainer (YIL-125 r6, 2026-08-06).

User asked for the literature slides WITHOUT abbreviations: full references,
the meaning of every column, and what checkmark / dash / text cells mean.
This builds a separate 4-slide file — the main deck is NOT touched:
  1. Setting comparison v2  (matrix + legend + [n] reference keys)
  2. What each column asks  (one full sentence per column)
  3. The works, one line each (no abbreviations)
  4. References             (verified against arXiv/publisher pages
                             2026-08-06; two entries marked venue-to-confirm)
Output: ppt/lit_review_explained.pptx
"""

import copy
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = f"{HERE}/ppt/lit_review_explained.pptx"
INK, INK2 = "0B0B0B", "52514E"
ORANGE, GREY, FILL = "EB6834", "8A8880", "F2F4F7"
GREEN = "1BAF7A"
EMU_IN = 914400


def IN(v):
    return Emu(int(v * EMU_IN))


def clone(prs, src):
    new = prs.slides.add_slide(src.slide_layout)
    for shp in list(new.shapes):
        shp._element.getparent().remove(shp._element)
    for shp in src.shapes:
        new.shapes._spTree.append(copy.deepcopy(shp._element))
    return new


def set_ph(slide, idx, text):
    for sh in slide.shapes:
        if sh.is_placeholder and sh.placeholder_format.idx == idx:
            p = sh.text_frame.paragraphs[0]
            if p.runs:
                p.runs[0].text = text
                for r in p.runs[1:]:
                    r._r.getparent().remove(r._r)
            else:
                p.text = text
            return sh
    return None


def bullets(slide, x, y, w, lines, size=14):
    tb = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(1))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (lvl, txt, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " if lvl == 0 else "   – ") + txt if lvl >= 0 else txt
        p.font.size = Pt(size if lvl <= 0 else size - 1.5)
        p.font.bold = bold
    return tb


def ntable(slide, x, y, w, col_ws, rows, font=9, header_font=9.5, row_h=0.32):
    nr, nc = len(rows), len(col_ws)
    gf = slide.shapes.add_table(nr, nc, IN(x), IN(y), IN(w), IN(row_h * nr))
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = True
    for j, cw in enumerate(col_ws):
        tbl.columns[j].width = IN(cw)
    for i, row in enumerate(rows):
        for j, cell in enumerate(row):
            txt, bold, color = (cell if isinstance(cell, tuple)
                                else (cell, False, INK))
            c = tbl.cell(i, j)
            c.margin_left = c.margin_right = IN(0.03)
            c.margin_top = c.margin_bottom = IN(0.01)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = c.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = txt
            p.font.size = Pt(header_font if i == 0 else font)
            p.font.bold = bold or (i == 0)
            p.font.color.rgb = RGBColor.from_string(color)
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
    return tbl


prs = Presentation(f"{HERE}/ppt/template.pptx")
n0 = len(prs.slides)
CONTENT = list(prs.slides)[3]

CK = ("✓", True, GREEN)
NO = ("—", False, GREY)
PART = ("partial", False, INK2)

# ------------------------------------------------- 1 matrix v2 with legend
s = clone(prs, CONTENT)
set_ph(s, 11, "Literature")
set_ph(s, 10, "Setting comparison")
ntable(s, 0.35, 1.42, 12.65,
       [3.30, 1.30, 1.42, 1.42, 1.28, 1.42, 1.28, 1.23],
       [
        [("work  [ref]", True, INK), "selective:\nchoose what\nto cover",
         "per-vehicle\nduties\n(o, d, t₀, B)", "arc-level, ONE\nfixed road\nnetwork",
         "time-dep.\ntravel times", "space–time\ncoverage\nobjective",
         "learned\namortised\nsolver", "unseen budget\nVALUES\n(zero-shot)"],
        ["[1] Kool et al. 2019 (neural Orienteering)", CK, NO, NO, NO, NO, CK, NO],
        ["[2] TOP-Former 2025 (neural Team Orienteering)", CK,
         ("budget only,\nhomogeneous", False, INK2), NO, NO, NO, CK, NO],
        ["[3] Lee & Ahn 2023 (multi-start TOP for UAS)", CK,
         ("start + fuel\nonly", False, INK2), NO, NO, NO, CK, NO],
        ["[4][5][6] MTPOMO · MVMoE · RouteFinder (multi-task VRP)", NO,
         ("dist-limit/open:\nvisit-all twins", False, INK2), NO, NO, NO, CK,
         ("variant-level,\nnot values", False, INK2)],
        ["[7] GOAL 2025 (generalist CO agent; OP task)", PART, NO, NO, NO, NO, CK, NO],
        ["[8] FM-MCVRP 2024 (fixed-graph supervised CVRP)", NO, NO,
         ("fixed graph,\nnode visits", False, INK2), NO, NO, CK, NO],
        ["[9] SED2AM 2025 (time-dependent delivery VRP)", NO, NO,
         ("real road\ntimes", False, INK2), CK, NO, CK, NO],
        ["[10] Neural CARP 2019–2026 (arc routing)", NO, NO, CK, NO, NO, CK, NO],
        ["[11][12][13] Drive-by-sensing OR (Han · Chen · Zhu)", CK, PART, CK,
         PART, CK, NO, NO],
        [("OURS (benchmark v3)", True, ORANGE), CK, CK, CK, CK, CK, CK, CK],
       ], font=8.4, header_font=8.6, row_h=0.445)
bullets(s, 2.05, 6.55, 10.7, [
    (0, "How to read a cell:  ✓ = the work has this property (as defined on the next slide) · "
        "— = it does not · grey text = it has only the PARTIAL form stated in the cell. "
        "[n] = reference list, last slide", True),
    (0, "No existing work is identical: every neighbour misses ≥ 2 columns. The bottom row of "
        "works has the right problem but no learning; the neural rows have learning but not "
        "this problem", False),
], size=10.5)

# ------------------------------------------------- 2 column key
s = clone(prs, CONTENT)
set_ph(s, 11, "Literature")
set_ph(s, 10, "What each column asks")
bullets(s, 0.7, 1.55, 12.1, [
    (0, "selective — choose what to cover:  may the fleet DECIDE which parts of the network to "
        "visit, maximising collected utility under a budget (orienteering / profits family)? "
        "A dash means visit-ALL: every customer or edge must be served at minimum cost "
        "(CVRP / CARP family) — a different optimisation problem", False),
    (0, "per-vehicle duties (o, d, t₀, B):  does EACH vehicle carry its own origin gate, "
        "destination gate, departure time AND budget? Partial forms: TOP-Former = one shared "
        "budget value, same depot for all; Lee & Ahn = own start + remaining fuel but no OD "
        "duties; the VRP foundation models = distance-limit / open-route attributes, which are "
        "the visit-all cousins of a budget and of o ≠ d", False),
    (0, "arc-level, ONE fixed road network:  is routing over the LINKS of one fixed network "
        "(revisits allowed, same graph every instance) rather than over freshly sampled points "
        "in the plane? Partial: FM-MCVRP fixes the graph but routes node visits; SED2AM uses "
        "real road travel times but is still node-based delivery", False),
    (0, "time-dependent travel times:  does the time to traverse a link depend on WHEN you "
        "enter it — c(i, t) — so the best route changes with the clock?", False),
    (0, "space–time coverage objective:  is the reward a UNION of (link, time) cells — where "
        "AND when you sensed — so re-sensing the same cell counts once, and sensing the same "
        "link at two different times counts twice?", False),
    (0, "learned amortised solver:  is there a trained neural policy that outputs a solution in "
        "milliseconds at deployment, amortising the cost of per-instance optimisation?", False),
    (0, "unseen budget VALUES (zero-shot):  is the model TESTED on budget values absent from "
        "every training label (continuous interpolation AND extrapolation), verified against "
        "an exact solver? 'variant-level, not values' = zero-shot to unseen COMBINATIONS of "
        "binary problem attributes only — no continuous physical dial", False),
], size=11.5)

# ------------------------------------------------- 3 works, one line each
s = clone(prs, CONTENT)
set_ph(s, 11, "Literature")
set_ph(s, 10, "The works, in one line each")
bullets(s, 0.7, 1.50, 12.1, [
    (0, "[1] Kool et al. 2019 — the founding attention model for learned routing; for the "
        "single-vehicle Orienteering Problem it introduced the budget-feasibility mask that "
        "our decoder generalises", False),
    (0, "[2] TOP-Former 2025 — centralised transformer that decodes a whole TEAM for the Team "
        "Orienteering Problem; closest neural relative (still: shared depot, one homogeneous "
        "budget, Euclidean points, static prizes)", False),
    (0, "[3] Lee & Ahn 2023 — Team Orienteering re-planning for aerial vehicles with "
        "heterogeneous current positions and remaining fuel; the heterogeneity precedent", False),
    (0, "[4] MTPOMO 2024 · [5] MVMoE 2024 · [6] RouteFinder 2025 — the multi-task VRP "
        "'foundation model' family: ONE network solves 16–48 visit-all delivery variants by "
        "composing binary attributes; our Proj([ρ, B/H]) conditioning is their recipe applied "
        "to a physical budget", False),
    (0, "[7] GOAL 2025 — one generalist backbone imitation-trained across many combinatorial "
        "problems (Orienteering among them); breadth, not our setting", False),
    (0, "[8] FM-MCVRP 2024 — supervised, LLM-style training on ONE fixed city graph for "
        "capacitated delivery; the closest TRAINING RECIPE to ours (fixed graph + supervised "
        "labels + joint fleet sequence)", False),
    (0, "[9] SED2AM 2025 — deep RL for multi-trip delivery with time-dependent travel times "
        "from real road data; the closest work on TIME DEPENDENCE (but visit-all delivery)", False),
    (0, "[10] Neural CARP 2019–2026 — the line of neural solvers for Capacitated ARC Routing "
        "(service on road edges); closest on ARC-LEVEL routing (but visit-all, static costs)", False),
    (0, "[11] Han et al. 2024 · [12] Chen, Qin & Sun 2024 · [13] Zhu et al. 2014 — the "
        "drive-by-sensing OR shelf: OUR problem statement (sensing coverage, budgets, real "
        "networks), solved by per-instance optimisation with NO learned solver — the gap we "
        "fill", False),
], size=11)

# ------------------------------------------------- 4 references
s = clone(prs, CONTENT)
set_ph(s, 11, "Literature")
set_ph(s, 10, "References")
bullets(s, 0.7, 1.42, 12.2, [
    (0, "[1] W. Kool, H. van Hoof, M. Welling. Attention, Learn to Solve Routing Problems! "
        "ICLR 2019. arXiv:1803.08475", False),
    (0, "[2] D. Fuertes, C. R. del-Blanco, F. Jaureguizar, N. García. TOP-Former: A Multi-Agent "
        "Transformer Approach for the Team Orienteering Problem. IEEE Trans. on Intelligent "
        "Transportation Systems, 2025. arXiv:2311.18662", False),
    (0, "[3] D. H. Lee, J. Ahn. Multi-Start Team Orienteering Problem for UAS Mission "
        "Re-Planning with Data-Efficient Deep Reinforcement Learning. Applied Intelligence, "
        "2024. arXiv:2303.01963", False),
    (0, "[4] F. Liu, X. Lin, Q. Zhang, X. Tong, M. Yuan. Multi-Task Learning for Routing "
        "Problem with Cross-Problem Zero-Shot Generalization. ACM SIGKDD 2024. "
        "arXiv:2402.16891  (“MTPOMO”)", False),
    (0, "[5] J. Zhou, Z. Cao, Y. Wu, W. Song, Y. Ma, J. Zhang, C. Xu. MVMoE: Multi-Task "
        "Vehicle Routing Solver with Mixture-of-Experts. ICML 2024. arXiv:2405.01029", False),
    (0, "[6] F. Berto et al. RouteFinder: Towards Foundation Models for Vehicle Routing "
        "Problems. TMLR, 2025. arXiv:2406.15007", False),
    (0, "[7] D. Drakulic, S. Michel, J.-M. Andreoli. GOAL: A Generalist Combinatorial "
        "Optimization Agent Learner. ICLR 2025. arXiv:2406.15079", False),
    (0, "[8] S. J. K. Chin, A. Srivastava, M. Winkenbach. Learning to Deliver: a Foundation "
        "Model for the Montreal Capacitated Vehicle Routing Problem. MIT, 2024. "
        "arXiv:2403.00026  (“FM-MCVRP”)", False),
    (0, "[9] A. Mozhdehi, Y. Wang, S. Sun, X. Wang. SED2AM: Solving Multi-Trip Time-Dependent "
        "Vehicle Routing Problem Using Deep Reinforcement Learning. ACM TKDD, 2025. "
        "arXiv:2503.04085", False),
    (0, "[10] Line of work, anchors: Learning to Solve Capacitated Arc Routing Problems by "
        "Policy Gradient (IEEE CEC 2019, doi:10.1109/CEC.2019.8790295) · A Neural Solver With "
        "Traversal-Based Feature Representation and Adjacent Attention for CARP (IEEE, 2025) · "
        "Direction-Aware Deep Policy Learning for Efficient Capacitated Arc Routing (Eng. "
        "Appl. AI, 2026)", False),
    (0, "[11] K. Han, W. Ji, Y. M. Nie, Z. Li, S. Liu. Exploring the Sensing Power of Mixed "
        "Vehicle Fleets. Transportation Research Part B 190:103066, 2024. "
        "doi:10.1016/j.trb.2024.103066", False),
    (0, "[12] X. Chen, G. Qin, J. Sun. Coordinated Routing Policy for Connected Vehicles to "
        "Monitor City-Wide Traffic. 2024. [venue to confirm]", False),
    (0, "[13] Zhu et al. Mobile Traffic Sensor Routing in Dynamic Transportation Systems. "
        "2014. [venue to confirm — IEEE ITS-family journal]", False),
    (0, "Also tracked (threats list): DeCoST, ICLR 2026, arXiv:2603.06260 — learning for "
        "Orienteering with time windows and variable profits", False),
], size=9.8)

xml = prs.slides._sldIdLst
for sld in list(xml)[:n0]:
    xml.remove(sld)
prs.save(OUT)
print(f"saved {OUT} with {len(prs.slides)} slides")
