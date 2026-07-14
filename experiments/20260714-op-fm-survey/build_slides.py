#!/usr/bin/env python3
"""Build the YIL-113 literature-survey deck by cloning slide 3 ("LR1"/MTPOMO) of
notes_discussion.pptx and swapping text, so header/subtitle/body formatting
(white banner header, big subtitle, Amasis MT Pro bulleted body) match exactly.

Reads the template for master/layout/theme + the slide-3 shapes to clone; writes a
NEW standalone .pptx (does NOT modify the template). Content = literature-reviewer
deliveries on YIL-114 (neural OP/TOP) and YIL-115 (VRP/CO foundation models).

Header ≤ ~24 chars and subtitle short (≈ 'LR1') so neither wraps over the body.
"""
import copy
from pptx import Presentation
from pptx.oxml.ns import qn

TEMPLATE = "/tmp/yil113_ppt/notes_discussion.pptx"
OUT = "/home/yilin/Research/Route_TSC_CART/experiments/20260714-op-fm-survey/op_fm_survey_slides.pptx"
SRC_IDX = 2  # slide 3 (0-based) = the LR1 style template

# ---- content DSL: (header, subtitle, [(kind, text)])
#   'H' framing/normal line  (level 0, bulleted like the original)
#   'L' bold section label    (level 0, bold)
#   'B' bullet item           (level 1, sub-bullet)
#   's' spacer                (blank line)
S = ("s", "")

SLIDES = [
 ("Literature Review (LR)", "Agenda", [
   ("H", "Scope (for YIL-113): how neural / foundation-model VRP solvers map to our robotaxi "
         "sensing-routing (team-orienteering) problem."),
   S,
   ("L", "Contents"),
   ("B", "Part A — Neural OP/TOP (5): Kool AM · POMO⚠ · TOP-Former★ · UAS-MSTOP · DTOP-SC⚠(OR)"),
   ("B", "Part B — VRP/CO foundation models (4): RouteFinder · FM-MCVRP★ · GOAL · UniCO"),
   ("B", "Synthesis — three unification philosophies (+ FM-MCVRP orthogonal)"),
   ("B", "Transfer — feasibility to route_cart_tsc; deltas = our novelty"),
   S,
   ("H", "Sources: full-text arXiv PDFs via literature-reviewer (YIL-114/115). ★ = closest to us, ⚠ = caveat."),
 ]),

 ("Neural OP/TOP (LR)", "A1 · Kool AM", [
   ("H", "Attention, Learn to Solve Routing Problems! (Kool et al., ICLR 2019) — founding neural method for the Orienteering Problem (OP)."),
   ("H", "OP = maximize Σ node prizes s.t. tour length ≤ budget T; single vehicle, depot→depot, visiting optional."),
   S,
   ("L", "Gap"),
   ("B", "Earlier learned routing (Pointer Nets + actor-critic) trained poorly; RNN encoders order-dependent & not parallel."),
   S,
   ("L", "Method"),
   ("B", "Transformer enc–dec, NO positional encoding (order-invariant); 3 layers, 8 heads, batch-norm."),
   ("B", "Subset selection: depot is a selectable action → choosing depot ENDS the route; visited-so-far = chosen subset."),
   ("B", "Budget feasibility = MASKING (not penalty/repair): mask node j when d(prev,j)+d(j,depot) > remaining T. ← transferable kernel."),
   ("B", "Train: REINFORCE + greedy-rollout baseline (no critic)."),
   S,
   ("L", "Result"),
   ("B", "OP n=20: AM sampling-1280 within 1.56% of Gurobi; SOTA GA (Compass) only ~2% better; beats OR-Tools/Tsili."),
 ]),

 ("Neural OP/TOP (LR)", "A2 · POMO  ⚠", [
   ("H", "POMO (Kwon et al., NeurIPS 2020) — a training/inference recipe on Kool's AM. ⚠ Does NOT solve OP/TOP (only TSP/CVRP/Knapsack)."),
   S,
   ("L", "Gap"),
   ("B", "AM's greedy-rollout baseline gives mostly-negative advantage; the first action over-dominates the solution."),
   S,
   ("L", "Method"),
   ("B", "Same AM backbone; removes the 'pick start node' step."),
   ("B", "Multiple optima: N fixed start nodes → N parallel rollouts; shared baseline b = mean of the N returns (no critic)."),
   ("B", "Inference: N multi-start greedy + ×8 instance augmentation."),
   S,
   ("L", "Relevance to us (honest)"),
   ("B", "LOW as a selective-routing method: multi-start symmetry does NOT hold for max-prize subset selection (different starts → different subsets)."),
   ("B", "Reusable parts only: AM + masking, shared-baseline REINFORCE, ×8 augmentation. (TSP100 gap 0.14%; no OP/TOP numbers.)"),
 ]),

 ("Neural OP/TOP (LR)", "A3 · TOP-Former  ★", [
   ("H", "TOP-Former (Fuertes et al., T-ITS 2025) — centralized multi-agent Transformer for the Team Orienteering Problem. ★ Closest to us."),
   ("H", "TOP = max team reward; per-vehicle budget T; shared depot; visiting optional; m = 2..5 vehicles."),
   S,
   ("L", "Gap"),
   ("B", "Prior multi-vehicle neural TOP is decentralized/sequential — each agent decides without seeing others' state → suboptimal."),
   S,
   ("L", "Method"),
   ("B", "Centralized enc–dec (no GNN); shared encoder run once; 3 blocks, dim 128, batch-norm."),
   ("B", "Simultaneous decoding: each step picks one node per agent; a chosen node is masked (−∞) for still-to-act agents."),
   ("B", "Per-vehicle budget via return-to-depot masking (Eq.18): mask nodes unreachable-and-return within remaining time t^a."),
   ("B", "Train: REINFORCE + greedy-rollout baseline (Kool-style, not POMO)."),
   S,
   ("L", "Result / deltas vs us"),
   ("B", "n100,m5: TOP-Former 82.79 (best, 0.00%) > ACO 81.39 (1.69%); Gurobi-60s only 61.9% gap. Inference ~4ms CPU / 0.22ms GPU."),
   ("B", "Deltas: single shared depot, homogeneous budget, STATIC reward, node-based, O(n²) scaling."),
 ]),

 ("Neural OP/TOP (LR)", "A4 · UAS-MSTOP", [
   ("H", "Multi-Start TOP for UAS re-planning (Lee & Ahn, KAIST, 2023) — TOP generalized to MID-MISSION re-planning."),
   ("H", "K vehicles start from different current positions with different remaining fuel f_k."),
   S,
   ("L", "Gap"),
   ("B", "AM/POMO assume all vehicles start at the depot (pre-planning); don't model heterogeneous starts + remaining energy."),
   S,
   ("L", "Method"),
   ("B", "Deep Dynamic Transformer (DDTM): 4 enc / 2 dec; RE-RUNS the encoder after each vehicle completes (graph state changes)."),
   ("B", "Nested loop: build one vehicle to return-depot, update instance, next vehicle from its own start."),
   ("B", "Budget = action masking on 'can't reach & return within remaining fuel'."),
   ("B", "Data-efficient training: instance-augmentation baseline (K=8 transforms) replaces greedy rollout → −30% epoch time."),
   S,
   ("L", "Result / relevance"),
   ("B", "MSTOP ×8N! gap: n10 0.19%, n20 0.78%, n≥50 reference-best; >90% instances gap=0 at n20."),
   ("B", "HIGH relevance if our fleet must re-plan mid-shift from current position / remaining charge."),
 ]),

 ("Neural OP/TOP (LR)", "A5 · DTOP-SC  ⚠", [
   ("H", "Dynamic TOP in Spatial Crowdsourcing (Wu et al., SCU, 2026) — workers=vehicles, tasks=profit nodes, tasks arrive online."),
   ("L", "⚠ Method class: pure OR / metaheuristic (NOT neural)"),
   ("B", "Scen-RH-ALNS: scenario-sampling rolling-horizon ALNS. No neural nets, no RL. Mis-listed as 'neural' in intake → reclassify as OR baseline."),
   S,
   ("L", "Method"),
   ("B", "Event-driven rolling horizon; each epoch solves a static HT-TOPTW snapshot with ALNS (destroy/repair + 2-opt + SA)."),
   ("B", "Scenario-sampling lookahead (Bent–Van Hentenryck consensus): sample S=15 futures + N_vir=5 virtual tasks, vote on first real task."),
   ("B", "Heterogeneous fleet: per-worker OD + time windows; selective (unrouted = not visited)."),
   S,
   ("L", "Result / relevance"),
   ("B", "vs MPA on 1161 instances (high dynamism): profit within 0.79–3.23%, but decision time 0.14s vs ~195s (2–3 orders faster)."),
   ("B", "Structurally very close to us; use as a strong non-learning / anticipatory-dispatch baseline."),
 ]),

 ("Routing Foundation(LR)", "LR2 · RouteFinder", [
   ("H", "RouteFinder (Berto et al., TMLR 09/2025) — a foundation model unifying 48 VRP variants. [Corrected: TMLR, not ICML.]"),
   ("H", "Unification = ATTRIBUTE COMPOSITION: every variant = a subset of one super-problem MDOVRPMBLTW."),
   S,
   ("L", "Key features"),
   ("B", "Unified VRP environment: any attribute subset (multi-depot/open/backhaul/distance-limit/time-windows; e.g. no TW → TW=[0,∞])."),
   ("B", "Global attribute embedding φ₀..φₖ injected into a DEEP encoder (vs MTPOMO/MVMoE: shallow decoder only)."),
   ("B", "Modern encoder (RF-TE): RMSNorm, pre-norm, SwiGLU, FlashAttention. + Mixed-Batch Training across variants."),
   ("B", "Efficient Adapter Layers (EAL): zero-pad projection W'=[W;0] for unseen attributes → few-shot fine-tune."),
   S,
   ("L", "Train / result"),
   ("B", "RL (REINFORCE + POMO shared baseline); A100 9–24h/model."),
   ("B", "Beats MTPOMO/MVMoE on all 48 variants; gap ~1–5% vs HGS-PyVRP/OR-Tools; inference 1–2s vs 10–20 min."),
 ]),

 ("Routing Foundation(LR)", "LR3 · FM-MCVRP  ★", [
   ("H", "FM-MCVRP (Chin, Winkenbach, Srivastava, MIT, 2024) — 'Learning to Deliver.' ★ Closest to our fixed-road-network daily instances."),
   ("H", "LLM-style SUPERVISED learning: a T5 enc-dec learns 'the next node to visit' from many sub-optimal historical solutions."),
   S,
   ("L", "Data setup (≈ our setting)"),
   ("B", "Fixed graph G′ = 10,001 nodes; each instance = a random node SUBSET (a day's demand). ALL instances = subgraphs of the SAME graph."),
   ("B", "38.1M instances (381 sizes × 100k), each labeled by ONE HGS run @5s (deliberately cheap/sub-optimal, mimics real historical data)."),
   ("B", "T5 (206M params), LM objective, curriculum learning (small→large), nucleus-sampling decode."),
   S,
   ("L", "Headline: SL BEATS its training data"),
   ("B", "vs its own HGS@5s labels, NS-1000 beats the labels at ≥100 customers (400 → −1.05%, statistically significant)."),
   ("B", "Mechanism: under a tight budget HGS degrades with scale faster than the model, which aggregates across millions of instances."),
   ("B", "Single model spans sizes 20–800 & all capacities; beats AM (diverges ≥400); within 3.02% of LKH-3 at 400."),
 ]),

 ("Routing Foundation(LR)", "LR4 · GOAL", [
   ("H", "GOAL (Drakulic, Michel, Andreoli, NAVER LABS, ICLR 2025) — a Generalist CO agent: one backbone + light per-task adapters."),
   ("H", "Unification = GENERALIST BACKBONE + ADAPTER: shared trunk learns commonality; adapters learn each problem's specifics."),
   S,
   ("L", "Key features"),
   ("B", "Shared codebook: project each problem's features to a small fixed rep, map through a shared codebook → related representations."),
   ("B", "Mixed-attention blocks: inject EDGE info into every attention kernel → arbitrary node/edge/instance features."),
   ("B", "Multi-type transformer: replicate blocks per node/edge type but SHARE parameters."),
   S,
   ("L", "Train / result"),
   ("B", "Imitation learning on expert trajectories; 8 CO tasks (ATSP/CVRP/CVRPTW/OP; JSSP/UMS; Knapsack; MVC). SL fine-tune > from-scratch."),
   ("B", "Single-task GOAL SOTA on 7/8 (except CVRP: 2.34% vs POMO 1.21% / RF-TE 1.50%); multi-task only slightly worse."),
 ]),

 ("Routing Foundation(LR)", "LR5 · UniCO", [
   ("H", "UniCO (Pan et al., SJTU Thinklab, ICLR 2025) — unify CO by PROBLEM REDUCTION to a matrix-encoded general TSP."),
   ("H", "Reduce every problem to one canonical form (general TSP over any positive cost matrix), solve once, transform the tour back."),
   S,
   ("L", "Key features"),
   ("B", "Reductions: ATSP, 2D-TSP, HCP (Hamiltonian cycle), 3SAT → all encoded as general-TSP cost matrices."),
   ("B", "MatPOENet (Graph Transformer + RL): MatNet + Pseudo One-hot Embedding (scale-free) + Mix-Score Attention."),
   ("B", "MatDIFFNet (Graph Diffusion + SL): discrete diffusion over tours, extends Euclidean-TSP diffusion to matrix TSP."),
   S,
   ("L", "Result"),
   ("B", "Unified MatPOENet*-8x: avg opt-gap ~1.4% at N≈20, high find-rate, beats LKH(10k/500) on some tasks."),
   ("B", "Cross-task transfer (3SAT): 50→200 fine-tune steps reach 95.96→97.08% find-rate (≈ from-scratch 96.08%) vs 17.92% no-tune."),
   S,
   ("H", "Note: the literature DB has TWO duplicate UniCO pages — flagged for human merge."),
 ]),

 ("Synthesis (LR)", "LR6 · unification", [
   ("H", "Three ways the VRP/CO foundation models 'unify' — plus one orthogonal axis closest to us."),
   S,
   ("L", "1. Attribute composition — RouteFinder"),
   ("B", "Variants = attribute subsets of one super-problem. Training: RL. Scope: 48 VRP variants."),
   ("L", "2. Generalist backbone + adapter — GOAL"),
   ("B", "Shared trunk learns commonality; adapters learn specifics. Imitation. 8 CO families. (RouteFinder EAL adjacent.)"),
   ("L", "3. Problem reduction — UniCO"),
   ("B", "Reduce everything to matrix general-TSP. RL/SL. ATSP/2D-TSP/HCP/3SAT."),
   S,
   ("L", "4. (Orthogonal) Fixed graph × LLM-style SL — FM-MCVRP"),
   ("B", "Not multi-problem unification: one problem on a fixed graph, learned from cheap historical solutions, beating them & generalizing across scale."),
   ("B", "→ Closest to our robotaxi FIXED road network."),
 ]),

 ("Transfer to our problem", "Transfer · OP/TOP", [
   ("H", "Ours: a fleet under per-vehicle budget SELECTS & SEQUENCES road SEGMENTS to MAXIMIZE sensing utility on a FIXED road network."),
   ("H", "Team-orienteering / VRP-with-profits (selective, max-utility) — NOT CVRP. MILP/Gurobi = ground truth & baseline to beat."),
   S,
   ("L", "Most transferable"),
   ("B", "TOP-Former (A3) & UAS-MSTOP (A4): learned MULTI-vehicle TOP — pick+sequence, max utility, per-vehicle return-to-depot budget masking."),
   ("B", "Kool AM (A1): reusable feasibility KERNEL — depot-as-terminal + 'visit only if you can still return within budget'."),
   ("B", "POMO (A2): training/inference parts only; its multi-start symmetry breaks for our selective objective."),
   ("B", "DTOP-SC (A5): OR baseline; scenario-sampling anticipation if utility is time-revealed."),
   S,
   ("L", "Our deltas (shared gap = our novelty)"),
   ("B", "All 5 are node-based, static-reward, mostly-homogeneous budget."),
   ("B", "Ours = arc/segment-level, real fixed network, TIME-VARYING utility → subset-selection over arcs + BUDGET masks (not capacity masks)."),
 ]),

 ("Transfer to our problem", "Transfer · FMs", [
   ("H", "Which foundation-model recipe transfers to our sensing-routing problem?"),
   S,
   ("L", "FM-MCVRP — most directly transferable (closest)"),
   ("B", "Our daily instances ARE subsets of one fixed road network = the FM-MCVRP premise. Recipe: SL on historical MILP/Gurobi solutions, sampled decoding."),
   ("B", "'SL beats sub-optimal labels' = strongest argument a fast solver trained on our Gurobi solutions can rival Gurobi under a time budget."),
   ("B", "Delta: swap visit-all min-cost for selective/max-utility → add budget masking + subset-selection decoding."),
   S,
   ("L", "RouteFinder / GOAL — adapter fine-tune path"),
   ("B", "Pretrain-then-adapt (EAL / GOAL adapters). Caution: these FMs unify VISIT-ALL MIN-COST; our selective/profit objective needs new representation + budget masks (not out-of-the-box)."),
   ("L", "UniCO — least direct"),
   ("B", "Reducing arc-level, selective, time-varying sensing to a matrix general-TSP is unnatural; keep as a conceptual alternative."),
   S,
   ("L", "Bottom line"),
   ("B", "FM-MCVRP-style SL on Gurobi labels over our fixed network + a TOP-Former/MSTOP-style decoder with budget masking = most promising. Prototype that first."),
 ]),
]


def clone_slide(prs, src):
    """Duplicate a text-only slide: copy every shape element under a new slide on the same layout."""
    new = prs.slides.add_slide(src.slide_layout)
    for shp in list(new.shapes):           # drop layout's auto-created placeholders
        shp._element.getparent().remove(shp._element)
    for shp in src.shapes:                 # deep-copy source shapes (header ph, subtitle ph, content box)
        new.shapes._spTree.append(copy.deepcopy(shp._element))
    return new


def para_first_run(p_el):
    return p_el.findall(qn("a:r"))[0]


def set_para_text(p_el, text, bold=False, size_pt=None):
    """Set a cloned paragraph's text via its first run; drop extra runs; keep run formatting."""
    runs = p_el.findall(qn("a:r"))
    if not runs:
        return
    t = runs[0].find(qn("a:t"))
    if t is None:
        t = runs[0].makeelement(qn("a:t"), {}); runs[0].append(t)
    t.text = text
    for r in runs[1:]:
        p_el.remove(r)
    if bold or size_pt:
        rPr = runs[0].find(qn("a:rPr"))
        if rPr is None:
            rPr = runs[0].makeelement(qn("a:rPr"), {}); runs[0].insert(0, rPr)
        if bold:
            rPr.set("b", "1")
        if size_pt:
            rPr.set("sz", str(int(size_pt * 100)))


def make_spacer(tmpl_l0):
    sp = copy.deepcopy(tmpl_l0)
    for r in sp.findall(qn("a:r"))[1:]:
        sp.remove(r)
    pPr = sp.find(qn("a:pPr"))
    if pPr is None:
        pPr = sp.makeelement(qn("a:pPr"), {}); sp.insert(0, pPr)
    for tag in ("a:buChar", "a:buAutoNum"):
        for e in pPr.findall(qn(tag)):
            pPr.remove(e)
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))
    set_para_text(sp, " ", size_pt=8)   # short blank line
    return sp


def make_title(tmpl_l0, text, size_pt=15):
    """A bold, un-bulleted full-paper-title line to head a per-paper slide."""
    tp = copy.deepcopy(tmpl_l0)
    for r in tp.findall(qn("a:r"))[1:]:
        tp.remove(r)
    pPr = tp.find(qn("a:pPr"))
    if pPr is None:
        pPr = tp.makeelement(qn("a:pPr"), {}); tp.insert(0, pPr)
    for tag in ("a:buChar", "a:buAutoNum"):
        for e in pPr.findall(qn(tag)):
            pPr.remove(e)
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))   # no bullet on the title line
    set_para_text(tp, text, bold=True, size_pt=size_pt)
    return tp


def content_box(slide):
    for sh in slide.shapes:
        if sh.name == "TextBox 4":
            return sh
    raise RuntimeError("content textbox not found")


def find_ph(slide, idx):
    for sh in slide.shapes:
        if sh.is_placeholder and sh.placeholder_format.idx == idx:
            return sh
    return None


# Full paper titles (verbatim from the literature-reviewer full-text reads on YIL-114/115),
# keyed by the slide tag (subtitle.split()[0]) so each per-paper slide is labelled with its
# actual article name — the abbreviations alone are hard to remember.
TITLES = {
 "A1":  "Attention, Learn to Solve Routing Problems!  (arXiv:1803.08475)",
 "A2":  "POMO: Policy Optimization with Multiple Optima for Reinforcement Learning  (arXiv:2010.16011)",
 "A3":  "TOP-Former: A Multi-Agent Transformer Approach for the Team Orienteering Problem  (arXiv:2311.18662)",
 "A4":  "Multi-Start Team Orienteering Problem for UAS Mission Re-Planning with Data-Efficient Deep RL  (arXiv:2303.01963)",
 "A5":  "The Dynamic Team Orienteering Problem in Spatial Crowdsourcing: A Scenario Sampling Approach  (arXiv:2601.11010)",
 "LR2": "RouteFinder: Towards Foundation Models for Vehicle Routing Problems  (arXiv:2406.15007)",
 "LR3": "Learning to Deliver: A Foundation Model for the Montreal Capacitated Vehicle Routing Problem  (arXiv:2403.00026)",
 "LR4": "GOAL: A Generalist Combinatorial Optimization Agent Learner  (arXiv:2406.15079)",
 "LR5": "On Unified Combinatorial Optimization via Problem Reduction to Matrix-Encoded General TSP  (OpenReview yEwakMNIex)",
}


def main():
    prs = Presentation(TEMPLATE)
    n_orig = len(prs.slides)
    src = prs.slides[SRC_IDX]

    # capture level-0 and level-1 paragraph templates from the source content box
    src_txbody = content_box(src).text_frame._txBody
    src_paras = src_txbody.findall(qn("a:p"))
    tmpl_l0 = tmpl_l1 = None
    for p in src_paras:
        lvl = p.find(qn("a:pPr")).get("lvl") if p.find(qn("a:pPr")) is not None else None
        has_run = p.find(qn("a:r")) is not None
        if has_run and (lvl in (None, "0")) and tmpl_l0 is None:
            tmpl_l0 = copy.deepcopy(p)
        if has_run and lvl == "1" and tmpl_l1 is None:
            tmpl_l1 = copy.deepcopy(p)
    assert tmpl_l0 is not None and tmpl_l1 is not None, "could not capture paragraph templates"

    for header, subtitle, paras in SLIDES:
        slide = clone_slide(prs, src)
        set_para_text(find_ph(slide, 11).text_frame.paragraphs[0]._p, header)   # white banner header
        set_para_text(find_ph(slide, 10).text_frame.paragraphs[0]._p, subtitle) # big subtitle (keeps 28pt)

        txbody = content_box(slide).text_frame._txBody
        for p in txbody.findall(qn("a:p")):                                     # clear cloned body paras
            txbody.remove(p)
        tag = subtitle.split()[0]                                              # label per-paper slides with the full title
        if tag in TITLES:
            txbody.append(make_title(tmpl_l0, TITLES[tag]))
        for kind, text in paras:
            if kind == "s":
                txbody.append(make_spacer(tmpl_l0)); continue
            tmpl = tmpl_l1 if kind == "B" else tmpl_l0
            newp = copy.deepcopy(tmpl)
            set_para_text(newp, text, bold=(kind == "L"), size_pt=(14 if kind == "B" else 15))
            txbody.append(newp)

    xml_slides = prs.slides._sldIdLst                                           # remove original template slides
    for sld in list(xml_slides)[:n_orig]:
        xml_slides.remove(sld)

    prs.save(OUT)
    print("saved", OUT, "with", len(prs.slides), "slides")


if __name__ == "__main__":
    main()
