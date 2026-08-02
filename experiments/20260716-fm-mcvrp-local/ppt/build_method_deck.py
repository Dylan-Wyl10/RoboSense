"""Build the methodology deck for YIL-113 from the user's template.

Clones template slides (keeps master/theme/banner), swaps header (ph idx 11)
and subtitle (ph idx 10), adds content as textboxes/pictures. Results slides
carry an explicit [pending] placeholder until extension-1 training completes.
"""
import copy

from pptx import Presentation
from pptx.util import Emu, Pt

BASE = "/home/yilin/Research/Route_TSC_CART/experiments/20260716-fm-mcvrp-local"
TPL = f"{BASE}/ppt/template.pptx"
OUT = f"{BASE}/ppt/method_deck_v1.pptx"
R = f"{BASE}/results"
INK = None  # keep template ink

W, H = 12192000, 6858000
EMU_IN = 914400


def clone(prs, src):
    new = prs.slides.add_slide(src.slide_layout)
    for shp in list(new.shapes):
        shp._element.getparent().remove(shp._element)
    for shp in src.shapes:
        new.shapes._spTree.append(copy.deepcopy(shp._element))
    return new


def set_ph(slide, idx, text):
    for sh in slide.shapes:
        if sh.is_placeholder and sh.placeholder_format.idx == idx:
            sh.text_frame.paragraphs[0].runs[0].text = text if \
                sh.text_frame.paragraphs[0].runs else None
            if not sh.text_frame.paragraphs[0].runs:
                sh.text_frame.paragraphs[0].text = text
            # drop extra runs
            for r in sh.text_frame.paragraphs[0].runs[1:]:
                r._r.getparent().remove(r._r)
            return sh
    return None


def add_bullets(slide, x_in, y_in, w_in, lines, size=14):
    tb = slide.shapes.add_textbox(Emu(int(x_in * EMU_IN)), Emu(int(y_in * EMU_IN)),
                                  Emu(int(w_in * EMU_IN)), Emu(int(1 * EMU_IN)))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (lvl, txt, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("• " if lvl == 0 else "   – ") + txt if lvl >= 0 else txt
        p.font.size = Pt(size if lvl <= 0 else size - 1.5)
        p.font.bold = bold
    return tb


def add_pic(slide, path, x_in, y_in, w_in):
    return slide.shapes.add_picture(path, Emu(int(x_in * EMU_IN)),
                                    Emu(int(y_in * EMU_IN)),
                                    width=Emu(int(w_in * EMU_IN)))


prs = Presentation(TPL)
n0 = len(prs.slides)
S = list(prs.slides)
title_s, agenda_s, content_s, thanks_s = S[0], S[1], S[3], S[4]

# 1 title
s = clone(prs, title_s)
set_ph(s, 0, "Yilin Wang    Aug 2026")
set_ph(s, 11, "Learning to Route for Sensing Coverage")
set_ph(s, 10, "A neural amortized solver trained on Gurobi labels — methodology")

# 2 agenda
s = clone(prs, agenda_s)
set_ph(s, 10, "202608")
set_ph(s, 11, "Agenda")
add_bullets(s, 1.2, 2.0, 10.5, [
    (0, "Problem: fleet sensing-coverage routing on a fixed network", False),
    (0, "Pipeline: offline label farm  →  supervised training  →  millisecond inference", False),
    (0, "Network & instance definition (bidirectional 4×4, 8 gates, 56 ODs)", False),
    (0, "Model architecture & the feasibility mask", False),
    (0, "Training recipe: data generation + teacher forcing", False),
    (0, "Evidence so far (3×3 / 4×4)  ·  extension-1 evaluation design [results pending]", False),
], size=17)

# 3 problem definition
s = clone(prs, content_s)
set_ph(s, 11, "Methodology")
set_ph(s, 10, "Problem")
add_bullets(s, 0.7, 1.65, 11.8, [
    (0, "One case (an 'operating day'): congestion δ ∈ {0,1}^80 per directed link · V vehicles · per-vehicle task (origin gate, destination gate, departure t0)", False),
    (0, "Vehicles run SIMULTANEOUSLY on one timeline; entering link i at time t occupies cells (i, t … t+c−1), travel time c(i,t) = (base(i)+t)//4 + 1 + δᵢ (street-symmetric base)", False),
    (0, "Objective (identical to the Gurobi MILP):   min  [ α₁·total_cost  −  α₂·fleet_coverage ] / (V·H),   α = 0.3 / 0.7", True),
    (0, "coverage = |union of occupied (link, time) cells| — overlap counts once; horizon H = 338 is the per-vehicle budget", False),
    (0, "Taxonomy: team-orienteering-type (selective, max-utility, budget) — NOT visit-all min-cost CVRP", False),
])

# 4 pipeline
s = clone(prs, content_s)
set_ph(s, 11, "Methodology")
set_ph(s, 10, "Pipeline")
add_pic(s, f"{R}/figA_pipeline.png", 0.9, 1.55, 11.4)
add_bullets(s, 0.7, 6.75, 11.8, [
    (0, "Gurobi = the only solution source (offline labels); the trained model replaces per-case solving at deployment — 2–4 ms, no solver calls", False),
])

# 5 network
s = clone(prs, content_s)
set_ph(s, 11, "Methodology")
set_ph(s, 10, "Network")
add_pic(s, f"{R}/fig8b_bigrid_optionb.png", 1.4, 1.55, 10.6)
add_bullets(s, 0.7, 6.65, 11.8, [
    (0, "Bidirectional 4×4: 80 directed links, no U-turns · 8 gates (4 corners + 4 midpoints) · all 56 ordered ODs feasible (heatmap) · horizon 338", False),
])

# 6 architecture
s = clone(prs, content_s)
set_ph(s, 11, "Methodology")
set_ph(s, 10, "Model")
add_pic(s, f"{R}/figB_architecture.png", 1.2, 1.5, 10.8)

# 7 feasibility mask
s = clone(prs, content_s)
set_ph(s, 11, "Methodology")
set_ph(s, 10, "Feasibility mask")
add_bullets(s, 0.7, 1.65, 11.8, [
    (0, "Three lines between logits and softmax:  logits[illegal] = −∞  →  illegal probability is EXACTLY zero (grammar-constrained decoding)", True),
    (1, "not a successor of the current link / U-turn  (road topology)", False),
    (1, "cannot reach THIS vehicle's destination gate within the horizon (time-dependent Dijkstra query)", False),
    (1, "SEP only when adjacent to own destination; EOS only after V SEPs", False),
    (0, "No parameters — rules guarantee feasibility; the network only ranks legal moves (learned from Gurobi)", False),
    (0, "Standard in the literature: Kool AM §3 (u=−∞; OP budget rule) · MTPOMO Eq.10 · FM-MCVRP §4.2 (softmax(G·Wo+M)) · TOP-Former Eq.18", False),
    (0, "Measured necessity: adjacency-only mask → 10/200 horizon-infeasible decodes; + budget mask → 200/200 feasible", False),
])

# 8 training recipe
s = clone(prs, content_s)
set_ph(s, 11, "Methodology")
set_ph(s, 10, "Training")
add_bullets(s, 0.7, 1.6, 11.8, [
    (0, "Data (offline label farm, resumable + watchdog-supervised):", True),
    (1, "δ ~ iid{0,1}^80 · V ∈ {2,3,4,6} · OD stratified by difficulty terciles (earliest-arrival matrix) · t0 ∈ {0..5}", False),
    (1, "labels: multi-commodity time-expanded flow MILP (Gurobi, MIPGap 2%, 60 s cap) — validated vs min-time bound, decomposition-consistent", False),
    (1, "5000 train + 500 test + 300 OD-zero-shot (4 held-out ODs) + 300 fleet-extrapolation (V ∈ {5,8})", False),
    (0, "Objective design note: coverage normalized by V·H with α 0.3/0.7 — otherwise optima collapse to min-time routes (no sensing behavior to learn)", False),
    (0, "Training: teacher forcing, next-link cross-entropy (PAD ignored); no solver and no generation inside the training loop; AdamW, batch 64", True),
    (0, "Model selection: per-epoch MASKED GREEDY generation on validation cases — select on true objective gap, not on CE loss", False),
    (0, "Data-volume policy: nested 1k/2k/5k slope experiment; extend the resumable farm only if the curve is still steep", False),
])

# 9 learned behavior (toy evidence)
s = clone(prs, content_s)
set_ph(s, 11, "Evidence (toy)")
set_ph(s, 10, "Learned coordination")
add_pic(s, f"{R}/fig2_cases.png", 0.4, 1.6, 6.9)
add_bullets(s, 7.5, 2.0, 5.6, [
    (0, "Decoder writes the fleet as ONE token sequence — later vehicles attend to earlier routes", False),
    (0, "Same prefix, different context: veh-2 flips 3:1.00 → 19:0.85 after seeing veh-1 occupy the 3-22-23-24 corridor (real trace)", True),
    (0, "'Don't duplicate coverage' is learned from joint-optimal labels, not hard-coded", False),
])

# 10 toy results
s = clone(prs, content_s)
set_ph(s, 11, "Evidence (toy)")
set_ph(s, 10, "vs Gurobi")
add_pic(s, f"{R}/fig6_pct_gaps.png", 0.6, 1.7, 12.0)
add_bullets(s, 0.7, 6.15, 11.8, [
    (0, "3×3: 55.5% exact-match, mean 0.26% · 4×4: 68.5%, mean 0.09% · 200/200 feasible · 2–4 ms/case (Gurobi: seconds) · 0.5 M params, <150 MB VRAM", False),
])

# 11 extension-1 eval design + placeholder
s = clone(prs, content_s)
set_ph(s, 11, "Extension 1")
set_ph(s, 10, "Evaluation")
add_bullets(s, 0.7, 1.65, 11.8, [
    (0, "Three-layer exam for the multi-OD / multi-fleet / staggered-departure model:", True),
    (1, "L1 same-distribution: 500 unseen cases (trained ODs, new δ / task mixes)", False),
    (1, "L2 OD ZERO-SHOT: 4 OD pairs never seen in training (G1→G5, G6→G2, G4→G8, G7→G3)", False),
    (1, "L3 fleet extrapolation: V ∈ {5, 8} (trained on {2,3,4,6})", False),
    (0, "Metrics: feasibility rate · % matching Gurobi · absolute + percentage objective gap · per-case CSV + route visualisations", False),
    (0, "[ RESULTS PENDING — label farm running (5000 train, ~3 h ETA at deck build time); slides to be filled after training ]", True),
])

# 12 summary / next
s = clone(prs, content_s)
set_ph(s, 11, "Summary")
set_ph(s, 10, "Next steps")
add_bullets(s, 0.7, 1.7, 11.8, [
    (0, "One recipe, verified twice and now scaling: Gurobi labels offline → small Transformer → guaranteed-feasible millisecond inference", True),
    (0, "Verification chain at every step: MILP ≡ brute force (small grids) · decomposition ≡ simulator · per-case CSVs vs Gurobi", False),
    (0, "Next: extension-1 results (3-seed) → data-scaling curve → 5×5 real network (760 CTM cells; topology extraction path verified, signals exactly replicable offline)", False),
    (0, "Foundation-model gap list: OD conditioning (done in ext-1) → wider instance distributions → coordinate-based encoding for cross-network transfer (paused by decision)", False),
])

xml = prs.slides._sldIdLst
for sld in list(xml)[:n0]:
    xml.remove(sld)
prs.save(OUT)
print("saved", OUT, "with", len(prs.slides), "slides")
